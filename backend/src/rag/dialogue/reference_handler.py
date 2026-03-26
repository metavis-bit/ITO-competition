"""
参考资料处理器 — 支持 PDF/Word/PPT/图片/视频上传与解析

对应 A04 要求:
  - 2c) 提供参考资料上传功能（支持PDF, Word, PPT, 图片, 视频等）
  - 3b) 对上传的参考资料进行内容解析（文本提取、视频关键帧分析或摘要生成）
"""

from __future__ import annotations

import os
from pathlib import Path
from typing import Any, Dict, List, Optional


class ReferenceHandler:
    """
    参考资料解析器。

    支持格式: PDF, Word(.docx), PPT(.pptx), 图片(png/jpg), 视频(mp4/avi)
    复用现有 parsers 模块的解析能力。

    Example:
        handler = ReferenceHandler(config_path="config.yaml")
        result = handler.parse_file("/path/to/doc.pdf", teacher_note="参照第3章的格式")
        print(result["text"][:500])
    """

    def __init__(self, config_path: str = "config.yaml"):
        self.config_path = config_path
        self._cfg = None

    @property
    def cfg(self) -> Dict[str, Any]:
        if self._cfg is None:
            from ..config import load_config
            self._cfg = load_config(self.config_path)
        return self._cfg

    def parse_file(
        self,
        file_path: str,
        teacher_note: str = "",
    ) -> Dict[str, Any]:
        """
        解析单个参考资料文件。

        Args:
            file_path: 文件路径
            teacher_note: 教师对此文件的说明

        Returns:
            {
                "file_name": str,
                "file_type": str,
                "text": str,           # 提取的文本内容
                "metadata": dict,      # 元信息
                "teacher_note": str,
                "parse_method": str,   # 使用的解析方法
            }
        """
        fp = Path(file_path)
        if not fp.exists():
            return {
                "file_name": fp.name,
                "file_type": "unknown",
                "text": f"文件不存在: {file_path}",
                "metadata": {},
                "teacher_note": teacher_note,
                "parse_method": "none",
            }

        ext = fp.suffix.lower()
        file_type = self._detect_type(ext)

        try:
            if file_type == "pdf":
                text, meta, method = self._parse_pdf(fp)
            elif file_type == "word":
                text, meta, method = self._parse_word(fp)
            elif file_type == "pptx":
                text, meta, method = self._parse_pptx(fp)
            elif file_type == "image":
                text, meta, method = self._parse_image(fp)
            elif file_type == "video":
                text, meta, method = self._parse_video(fp)
            elif file_type == "text":
                text = fp.read_text(encoding="utf-8", errors="ignore")
                meta = {"chars": len(text)}
                method = "direct_read"
            else:
                text = f"不支持的文件格式: {ext}"
                meta = {}
                method = "unsupported"
        except Exception as e:
            text = f"解析失败: {e}"
            meta = {"error": str(e)}
            method = "error"

        return {
            "file_name": fp.name,
            "file_type": file_type,
            "text": text,
            "metadata": meta,
            "teacher_note": teacher_note,
            "parse_method": method,
        }

    def parse_files(
        self,
        file_paths: List[str],
        teacher_notes: Optional[List[str]] = None,
    ) -> List[Dict[str, Any]]:
        """批量解析多个参考资料"""
        notes = teacher_notes or [""] * len(file_paths)
        return [
            self.parse_file(fp, note)
            for fp, note in zip(file_paths, notes)
        ]

    def get_combined_text(self, results: List[Dict[str, Any]], max_chars: int = 10000) -> str:
        """将多个解析结果合并为一段文本（用于LLM输入）"""
        parts = []
        total = 0
        for r in results:
            header = f"[参考资料: {r['file_name']} ({r['file_type']})]"
            if r.get("teacher_note"):
                header += f"\n教师说明: {r['teacher_note']}"
            text = r.get("text", "")[:max_chars - total]
            parts.append(f"{header}\n{text}")
            total += len(text) + len(header)
            if total >= max_chars:
                break
        return "\n\n---\n\n".join(parts)

    # ─────────────────────────────────────────
    # 内部解析方法
    # ─────────────────────────────────────────

    @staticmethod
    def _detect_type(ext: str) -> str:
        mapping = {
            ".pdf": "pdf",
            ".doc": "word", ".docx": "word",
            ".ppt": "pptx", ".pptx": "pptx",
            ".png": "image", ".jpg": "image", ".jpeg": "image",
            ".gif": "image", ".bmp": "image", ".webp": "image",
            ".mp4": "video", ".avi": "video", ".mkv": "video",
            ".mov": "video", ".flv": "video",
            ".mp3": "audio", ".wav": "audio",
            ".txt": "text", ".md": "text", ".csv": "text",
        }
        return mapping.get(ext, "unknown")

    def _parse_pdf(self, fp: Path):
        """PDF解析 — 优先用 pdfplumber，退化到 PyPDF2"""
        try:
            import pdfplumber
            texts = []
            with pdfplumber.open(str(fp)) as pdf:
                for i, page in enumerate(pdf.pages):
                    t = page.extract_text() or ""
                    if t.strip():
                        texts.append(f"[第{i+1}页]\n{t}")
            text = "\n\n".join(texts)
            meta = {"pages": len(texts)}
            return text, meta, "pdfplumber"
        except ImportError:
            pass

        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(str(fp))
            texts = []
            for i, page in enumerate(reader.pages):
                t = page.extract_text() or ""
                if t.strip():
                    texts.append(f"[第{i+1}页]\n{t}")
            return "\n\n".join(texts), {"pages": len(texts)}, "PyPDF2"
        except ImportError:
            return "需要安装 pdfplumber 或 PyPDF2", {}, "missing_lib"

    def _parse_word(self, fp: Path):
        """Word解析"""
        try:
            import docx
            doc = docx.Document(str(fp))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            text = "\n".join(paragraphs)
            meta = {"paragraphs": len(paragraphs)}
            return text, meta, "python-docx"
        except ImportError:
            return "需要安装 python-docx", {}, "missing_lib"

    def _parse_pptx(self, fp: Path):
        """PPT解析"""
        try:
            from pptx import Presentation
            prs = Presentation(str(fp))
            texts = []
            for i, slide in enumerate(prs.slides):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text)
                if slide_texts:
                    texts.append(f"[幻灯片{i+1}]\n" + "\n".join(slide_texts))
            text = "\n\n".join(texts)
            meta = {"slides": len(prs.slides)}
            return text, meta, "python-pptx"
        except ImportError:
            return "需要安装 python-pptx", {}, "missing_lib"

    def _parse_image(self, fp: Path):
        """图片解析 — 优先 OCR，必要时使用 VLM 补充语义描述"""
        # 1) OCR
        try:
            from ..tools.ocr import ocr_image

            ocr_result = ocr_image(str(fp))
            text = (ocr_result.text or "").strip()
            if text and len(text) > 20:
                return (
                    text,
                    {
                        "method": "ocr",
                        "engine": ocr_result.engine,
                        "mean_conf": ocr_result.mean_conf,
                    },
                    "ocr",
                )
        except Exception:
            pass

        # 2) VLM
        try:
            from ..tools.vlm import describe_image

            vlm_cfg = self.cfg.get("vlm", {}) or {}
            if vlm_cfg.get("enabled"):
                vr = describe_image(
                    str(fp),
                    "请描述这张高中生物教学图片中的关键知识点、图示结构与可用于课堂讲解的信息。",
                    vlm_cfg,
                )
                if vr and (vr.text or "").strip():
                    return vr.text.strip(), {"method": "vlm", "engine": vr.engine}, "vlm"
        except Exception:
            pass

        return f"图片文件: {fp.name}（OCR/VLM未就绪，建议安装OCR或配置视觉模型）", {}, "placeholder"
    def _parse_video(self, fp: Path):
        """
        视频解析 — 音频转写 + 关键帧抽取 + VLM/OCR摘要

        对应 A04 3b): 视频关键帧分析或摘要生成
        """
        parts = []
        meta = {}

        # 1) ASR 语音转文字（视频先抽音频）
        try:
            from ..tools.asr import extract_audio_with_ffmpeg, transcribe_audio
            from ..tools.media import has_ffmpeg
            import tempfile

            if has_ffmpeg():
                with tempfile.TemporaryDirectory() as tmp_dir:
                    wav_path = Path(tmp_dir) / "audio.wav"
                    extract_audio_with_ffmpeg(str(fp), str(wav_path))
                    asr_result = transcribe_audio(str(wav_path))
                    asr_text = (asr_result.text or "").strip()
                    if asr_text:
                        parts.append(f"[语音转写]\n{asr_text}")
                        meta["asr"] = True
                        meta["asr_engine"] = asr_result.engine
            else:
                meta["asr_error"] = "ffmpeg unavailable"
        except Exception as e:
            meta["asr_error"] = str(e)

        # 2) 关键帧提取 + VLM/OCR 描述
        try:
            from ..tools.media import extract_frames_with_ffmpeg, has_ffmpeg
            from ..tools.vlm import describe_image
            from ..tools.ocr import ocr_image
            import tempfile

            if has_ffmpeg():
                with tempfile.TemporaryDirectory() as tmp_dir:
                    frames = extract_frames_with_ffmpeg(
                        str(fp),
                        out_dir=tmp_dir,
                        interval_sec=10,
                        max_frames=5,
                    )
                    if frames:
                        meta["frames_extracted"] = len(frames)
                        vlm_cfg = self.cfg.get("vlm", {}) or {}
                        for i, frame_path in enumerate(frames[:3]):
                            frame_text = ""
                            if vlm_cfg.get("enabled"):
                                try:
                                    vr = describe_image(
                                        str(frame_path),
                                        "请概述该关键帧中与高中生物教学相关的核心信息、图示关系与术语。",
                                        vlm_cfg,
                                    )
                                    if vr and (vr.text or "").strip():
                                        frame_text = vr.text.strip()
                                        meta["vlm_frames"] = True
                                except Exception:
                                    frame_text = ""

                            if not frame_text:
                                try:
                                    ocr_result = ocr_image(str(frame_path))
                                    frame_text = (ocr_result.text or "").strip()
                                except Exception:
                                    frame_text = ""

                            if frame_text:
                                parts.append(f"[关键帧{i + 1}]\n{frame_text}")
            else:
                meta["frame_error"] = "ffmpeg unavailable"
        except Exception as e:
            meta["frame_error"] = str(e)

        if parts:
            method = "asr" + (
                "+vlm_frames"
                if meta.get("vlm_frames")
                else "+frames"
                if meta.get("frames_extracted")
                else ""
            )
            return "\n\n".join(parts), meta, method

        return (
            f"视频文件: {fp.name}（ASR和关键帧提取未就绪，需要ffmpeg及ASR/OCR/VLM依赖）",
            meta,
            "placeholder",
        )

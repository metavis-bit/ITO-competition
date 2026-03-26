"""
PPT预览工具 — 将PPTX转为缩略图用于Gradio展示

使用 LibreOffice 将 pptx → pdf → png 图片
"""

from __future__ import annotations

import subprocess
import tempfile
from pathlib import Path
from typing import List, Optional


def pptx_to_images(pptx_path: str, max_pages: int = 20, dpi: int = 150) -> List[str]:
    """
    将PPTX转为PNG图片列表。

    流程: pptx → (LibreOffice) → pdf → (pdf2image) → png list

    Args:
        pptx_path: PPTX文件路径
        max_pages: 最大转换页数
        dpi: 图片分辨率

    Returns:
        PNG文件路径列表
    """
    pptx_path = str(pptx_path)
    if not Path(pptx_path).exists():
        return []

    with tempfile.TemporaryDirectory() as tmp_dir:
        # Step 1: pptx → pdf (via LibreOffice)
        try:
            subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "pdf",
                 "--outdir", tmp_dir, pptx_path],
                capture_output=True, timeout=60,
            )
        except (subprocess.TimeoutExpired, FileNotFoundError):
            return []

        pdf_files = list(Path(tmp_dir).glob("*.pdf"))
        if not pdf_files:
            return []
        pdf_path = str(pdf_files[0])

        # Step 2: pdf → png images
        try:
            from pdf2image import convert_from_path
            images = convert_from_path(pdf_path, dpi=dpi, first_page=1,
                                       last_page=max_pages)
        except Exception:
            return []

        # Step 3: 保存到持久目录
        out_dir = Path("outputs/ppt_preview")
        out_dir.mkdir(parents=True, exist_ok=True)

        stem = Path(pptx_path).stem
        paths = []
        for i, img in enumerate(images):
            out_path = str(out_dir / f"{stem}_slide_{i+1}.png")
            img.save(out_path, "PNG")
            paths.append(out_path)

        return paths


def pptx_preview_html(pptx_path: str, max_pages: int = 15) -> str:
    """
    生成PPT预览的HTML画廊（嵌入Gradio HTML组件）。
    """
    images = pptx_to_images(pptx_path, max_pages=max_pages, dpi=120)

    if not images:
        # fallback: 文本预览
        return _text_fallback(pptx_path)

    html_parts = ['<div style="display:flex;flex-direction:column;gap:12px;padding:8px;">']
    for i, img_path in enumerate(images):
        # 读取图片为base64内嵌
        import base64
        with open(img_path, "rb") as f:
            b64 = base64.b64encode(f.read()).decode()
        html_parts.append(
            f'<div style="border:1px solid #ddd;border-radius:8px;overflow:hidden;">'
            f'<div style="background:#f0f0f0;padding:4px 12px;font-size:12px;color:#666;">第 {i+1} 页</div>'
            f'<img src="data:image/png;base64,{b64}" style="width:100%;display:block;" />'
            f'</div>'
        )
    html_parts.append('</div>')
    return "\n".join(html_parts)


def _text_fallback(pptx_path: str) -> str:
    """LibreOffice不可用时的文本预览回退"""
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        lines = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip()[:60])
            content = " | ".join(texts[:3]) if texts else "(视觉元素)"
            lines.append(f"<p><b>📄 第{i}页:</b> {content}</p>")
        return "<div style='padding:8px;'>" + "\n".join(lines) + "</div>"
    except Exception as e:
        return f"<p>预览失败: {e}</p>"

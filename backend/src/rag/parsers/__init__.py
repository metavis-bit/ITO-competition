from __future__ import annotations

import hashlib
import os
import re
import subprocess
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Dict, List, Optional

from ..tools.asr import extract_audio_with_ffmpeg, save_asr_result, transcribe_audio
from ..tools.media import convert_audio_to_wav, extract_frames_with_ffmpeg, has_ffmpeg
from ..tools.ocr import ocr_image, save_ocr_result
from ..tools.vlm import describe_image, save_vlm_result


@dataclass
class ParsedDocument:
    text: str
    meta: Dict[str, Any]


def _sha1(s: str) -> str:
    return hashlib.sha1(s.encode("utf-8")).hexdigest()


def _safe_read_text(path: Path) -> str:
    for enc in ["utf-8", "utf-8-sig", "gbk", "latin-1"]:
        try:
            return path.read_text(encoding=enc)
        except Exception:
            continue
    return ""


def _relpath(path: Path, root_dir: str) -> str:
    try:
        return str(path.relative_to(Path(root_dir)))
    except Exception:
        return path.name


def _libreoffice_convert(in_path: Path, out_dir: Path, to_ext: str) -> Optional[Path]:
    """Convert office formats via LibreOffice (soffice).

    Returns converted file path or None.
    """
    from shutil import which

    if which("soffice") is None:
        return None

    out_dir.mkdir(parents=True, exist_ok=True)
    cmd = [
        "soffice",
        "--headless",
        "--nologo",
        "--nofirststartwizard",
        "--convert-to",
        to_ext,
        "--outdir",
        str(out_dir),
        str(in_path),
    ]
    try:
        subprocess.run(cmd, check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    except Exception:
        return None

    # LibreOffice keeps original stem
    cand = out_dir / f"{in_path.stem}.{to_ext}"
    if cand.exists():
        return cand
    # Sometimes it may output uppercase or different
    for p in out_dir.glob(in_path.stem + ".*"):
        if p.suffix.lower().lstrip(".") == to_ext.lower():
            return p
    return None


def _table_to_markdown(table: List[List[Any]]) -> str:
    """Convert a 2D table (list-of-rows) to a Markdown table.

    We keep it dependency-free (no pandas/tabulate).
    """
    if not table:
        return ""

    # Normalize
    rows: List[List[str]] = []
    max_cols = 0
    for r in table:
        if r is None:
            continue
        row = ["" if c is None else str(c).replace("\n", " ").strip() for c in (r or [])]
        rows.append(row)
        max_cols = max(max_cols, len(row))
    if not rows or max_cols == 0:
        return ""

    for r in rows:
        if len(r) < max_cols:
            r.extend([""] * (max_cols - len(r)))

    header = rows[0]
    body = rows[1:] if len(rows) > 1 else []

    def _line(cells: List[str]) -> str:
        return "| " + " | ".join([c.replace("|", "\\|") for c in cells]) + " |"

    md = [_line(header), "| " + " | ".join(["---"] * max_cols) + " |"]
    for r in body:
        md.append(_line(r))
    return "\n".join(md).strip()


def _clean_lines(text: str, *, regexes: List[str], drop_lines: Optional[set[str]] = None) -> str:
    """Line-wise cleaning for PDF/HTML extracted text."""
    drop_lines = drop_lines or set()
    out_lines: List[str] = []
    for raw in (text or "").splitlines():
        line = raw.strip()
        if not line:
            continue
        if line in drop_lines:
            continue
        matched = False
        for rgx in regexes:
            try:
                if re.match(rgx, line):
                    matched = True
                    break
            except Exception:
                continue
        if matched:
            continue
        out_lines.append(line)
    return "\n".join(out_lines).strip()


def _find_repeated_header_footer(pages: List[str], ratio: float = 0.6) -> set[str]:
    """Detect repeated header/footer lines across pages.

    Heuristic: take first and last non-empty line of each page.
    If a line repeats on >= ratio of pages, drop it.
    """

    from collections import Counter

    first_lines: List[str] = []
    last_lines: List[str] = []
    for t in pages:
        lines = [x.strip() for x in (t or "").splitlines() if x and x.strip()]
        if not lines:
            continue
        first_lines.append(lines[0])
        last_lines.append(lines[-1])

    if not first_lines and not last_lines:
        return set()

    n = max(1, len(pages))
    thr = max(2, int(n * float(ratio)))
    c = Counter(first_lines + last_lines)
    return {k for k, v in c.items() if v >= thr}


def _is_chart_like(image_path: str, line_threshold: int = 25) -> bool:
    """Best-effort chart/diagram detection.

    If OpenCV is installed, use HoughLinesP count as a rough signal of
    "many straight lines" (common in charts/flowcharts).

    Fail-soft: returns False when OpenCV is unavailable.
    """
    try:
        import cv2  # type: ignore
        import numpy as np  # type: ignore
    except Exception:
        return False

    try:
        img = cv2.imread(image_path)
        if img is None:
            return False
        gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
        edges = cv2.Canny(gray, 50, 150, apertureSize=3)
        lines = cv2.HoughLinesP(edges, 1, np.pi / 180, threshold=80, minLineLength=30, maxLineGap=5)
        cnt = 0 if lines is None else int(len(lines))
        return cnt >= int(line_threshold)
    except Exception:
        return False


def parse_file(
    file_path: str,
    *,
    root_dir: str,
    source_type: str,
    session_id: str,
    assets_dir: Optional[str],
    parse_cfg: Dict[str, Any],
) -> List[ParsedDocument]:
    """Parse a file into *retrievable evidence*.

    This is a pragmatic multi-format parser for Edu-RAG.

    Supported (best-effort):
      - PDF: text + page images (optional OCR fallback)
      - DOCX / DOC: paragraphs + tables
      - PPTX / PPT: slide texts (+notes)
      - XLSX / XLS: per-sheet table text
      - Image: OCR text
      - Audio: ASR transcript (segmented)
      - Video: ASR transcript + optional frame OCR
      - TXT/MD: raw text
      - HTML: main text extraction

    Returns a list of ParsedDocument, each with text + rich meta pointers.
    """
    p = Path(file_path)
    if not p.exists() or not p.is_file():
        return []

    rel = _relpath(p, root_dir)
    ext = p.suffix.lower()

    # Assets layout: assets_dir/<session_id>/<file_id>/...
    assets_root = Path(assets_dir) if assets_dir else None
    file_id = _sha1(f"{source_type}:{session_id}:{rel}")

    base_meta = {
        "source_type": source_type,
        "session_id": session_id,
        "abs_path": str(p.resolve()),
        "rel_path": rel,
        "source": rel,
        "file_ext": ext,
        "file_id": file_id,
    }

    # ---------- format dispatch ----------
    if ext in [".txt", ".md", ".csv", ".log"]:
        txt = _safe_read_text(p).strip()
        if not txt:
            return []
        meta = dict(base_meta)
        meta.update({"loc": rel, "part": "text"})
        return [ParsedDocument(text=txt, meta=meta)]

    if ext in [".html", ".htm"]:
        try:
            from bs4 import BeautifulSoup  # type: ignore
        except Exception:
            txt = _safe_read_text(p)
            meta = dict(base_meta)
            meta.update({"loc": rel, "part": "html"})
            return [ParsedDocument(text=txt.strip(), meta=meta)]

        html = _safe_read_text(p)
        soup = BeautifulSoup(html, "lxml")
        # remove noise
        for tag in soup(["script", "style", "noscript"]):
            try:
                tag.decompose()
            except Exception:
                pass
        txt = soup.get_text("\n")
        txt = "\n".join([x.strip() for x in txt.splitlines() if x.strip()])
        meta = dict(base_meta)
        meta.update({"loc": rel, "part": "html"})
        return [ParsedDocument(text=txt.strip(), meta=meta)]

    if ext in [".pdf"]:
        """PDF parsing strategy (aligned with队长分工文档, but keep our stronger fallbacks):

        1) Prefer pdfplumber for text + tables extraction.
        2) Clean headers/footers/page numbers.
        3) If a page has too few characters (e.g. < 50), treat it as scanned and run OCR.
        4) (Optional) render page images for evidence preview/citation.
        """

        try:
            import fitz  # PyMuPDF for rendering
        except Exception as e:
            raise RuntimeError("PyMuPDF is required for PDF parsing: pip install pymupdf") from e

        pdf_cfg = (parse_cfg.get("pdf") or {})
        use_pdfplumber = bool(pdf_cfg.get("use_pdfplumber", True))
        extract_tables = bool(pdf_cfg.get("extract_tables", True))
        render_pages = bool(pdf_cfg.get("render_page_images", True))
        ocr_fallback = bool(pdf_cfg.get("ocr_if_no_text", True))
        ocr_threshold = int(pdf_cfg.get("ocr_threshold_chars", 50))
        clean_regexes = list(pdf_cfg.get("clean_regexes") or [])
        remove_rep = bool(pdf_cfg.get("remove_repeated_header_footer", True))
        rep_ratio = float(pdf_cfg.get("repeated_line_ratio", 0.6))

        # ---- pass 1: extract raw page texts + tables (no OCR) ----
        raw_pages: List[str] = []
        page_tables_md: List[List[str]] = []

        if use_pdfplumber:
            try:
                import pdfplumber  # type: ignore
            except Exception:
                pdfplumber = None
        else:
            pdfplumber = None

        if pdfplumber is not None:
            with pdfplumber.open(str(p)) as pdf:
                for pg in pdf.pages:
                    t = (pg.extract_text() or "").strip()
                    tabs_md: List[str] = []
                    if extract_tables:
                        try:
                            for tb in (pg.extract_tables() or []):
                                md = _table_to_markdown(tb or [])
                                if md:
                                    tabs_md.append(md)
                        except Exception:
                            tabs_md = []
                    raw_pages.append(t)
                    page_tables_md.append(tabs_md)
        else:
            # fallback: PyMuPDF text
            doc0 = fitz.open(str(p))
            for i in range(len(doc0)):
                page = doc0.load_page(i)
                raw_pages.append((page.get_text("text") or "").strip())
                page_tables_md.append([])
            doc0.close()

        drop_lines = _find_repeated_header_footer(raw_pages, ratio=rep_ratio) if remove_rep else set()

        # ---- pass 2: build ParsedDocument per page (with OCR + page images) ----
        docs: List[ParsedDocument] = []
        doc = fitz.open(str(p))
        for i in range(len(doc)):
            page = doc.load_page(i)
            page_no = i + 1

            raw_text = raw_pages[i] if i < len(raw_pages) else (page.get_text("text") or "")
            cleaned = _clean_lines(raw_text or "", regexes=clean_regexes, drop_lines=drop_lines)

            # attach tables markdown
            tabs_md = page_tables_md[i] if i < len(page_tables_md) else []
            if tabs_md:
                if cleaned:
                    cleaned = cleaned + "\n\n[表格]" + "\n\n" + "\n\n".join(tabs_md)
                else:
                    cleaned = "[表格]\n\n" + "\n\n".join(tabs_md)

            meta = dict(base_meta)
            meta.update({
                "page": page_no,
                "part": f"page_{page_no}",
                "loc": f"{rel}#page={page_no}",
                "tables": len(tabs_md),
            })

            page_img_path: Optional[Path] = None
            if assets_root and render_pages:
                out_dir = assets_root / "pdf" / file_id
                out_dir.mkdir(parents=True, exist_ok=True)
                page_img_path = out_dir / f"page_{page_no:04d}.png"
                if not page_img_path.exists():
                    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                    pix.save(str(page_img_path))
                meta["page_image_path"] = str(page_img_path)

            # OCR fallback for scanned pages
            text_for_store = cleaned
            if ocr_fallback and page_img_path and page_img_path.exists():
                # either empty, or too short (e.g., scanned)
                if (not cleaned) or (len(cleaned) < int(ocr_threshold)):
                    try:
                        o = ocr_image(str(page_img_path))
                        ocr_txt = (o.text or "").strip()
                        if ocr_txt:
                            text_for_store = ocr_txt
                        ocr_path = page_img_path.with_suffix(".ocr.json")
                        save_ocr_result(o, str(ocr_path))
                        meta.update({"ocr_path": str(ocr_path), "ocr_engine": o.engine, "ocr_mean_conf": o.mean_conf})
                    except Exception:
                        pass

            if text_for_store:
                docs.append(ParsedDocument(text=text_for_store, meta=meta))

        try:
            doc.close()
        except Exception:
            pass
        return docs

    if ext in [".docx"]:
        try:
            from docx import Document  # type: ignore
        except Exception as e:
            raise RuntimeError("python-docx is required: pip install python-docx") from e

        doc = Document(str(p))
        parts: List[str] = []
        for para in doc.paragraphs:
            t = (para.text or "").strip()
            if t:
                parts.append(t)
        # tables
        for table in getattr(doc, "tables", []) or []:
            for row in table.rows:
                cells = [(c.text or "").strip() for c in row.cells]
                cells = [c for c in cells if c]
                if cells:
                    parts.append("\t".join(cells))
        txt = "\n".join(parts).strip()
        if not txt:
            return []
        meta = dict(base_meta)
        meta.update({"loc": rel, "part": "docx"})
        return [ParsedDocument(text=txt, meta=meta)]

    if ext in [".doc"]:
        # Convert to docx if possible (keep original file pointers in meta)
        out_dir = (assets_root / "converted" / file_id) if assets_root else (Path(p.parent) / ".converted")
        converted = _libreoffice_convert(p, out_dir, "docx")
        if converted:
            docs = parse_file(
                str(converted),
                root_dir=str(converted.parent),
                source_type=source_type,
                session_id=session_id,
                assets_dir=assets_dir,
                parse_cfg=parse_cfg,
            )
            for d in docs:
                d.meta["orig_abs_path"] = d.meta.get("abs_path")
                d.meta["orig_rel_path"] = d.meta.get("rel_path")
                d.meta["abs_path"] = str(p.resolve())
                d.meta["rel_path"] = rel
                d.meta["source"] = rel
                d.meta["loc"] = rel
                d.meta["file_ext"] = ext
            return docs

        # fallback: try antiword
        from shutil import which

        if which("antiword") is not None:
            try:
                res = subprocess.run(["antiword", str(p)], check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
                txt = res.stdout.decode("utf-8", errors="ignore").strip()
                if txt:
                    meta = dict(base_meta)
                    meta.update({"loc": rel, "part": "doc"})
                    return [ParsedDocument(text=txt, meta=meta)]
            except Exception:
                pass

        # If we can't parse, return empty
        return []

    if ext in [".pptx"]:
        try:
            from pptx import Presentation  # type: ignore
        except Exception as e:
            raise RuntimeError("python-pptx is required: pip install python-pptx") from e

        prs = Presentation(str(p))
        docs: List[ParsedDocument] = []
        for i, slide in enumerate(prs.slides, start=1):
            texts: List[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    t = (shape.text or "").strip()
                    if t:
                        texts.append(t)
            # slide notes
            try:
                if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
                    nt = (slide.notes_slide.notes_text_frame.text or "").strip()
                    if nt:
                        texts.append("[Notes] " + nt)
            except Exception:
                pass

            txt = "\n".join(texts).strip()
            if not txt:
                continue
            meta = dict(base_meta)
            meta.update({
                "slide": i,
                "part": f"slide_{i}",
                "loc": f"{rel}#slide={i}",
            })
            docs.append(ParsedDocument(text=txt, meta=meta))
        return docs

    if ext in [".ppt"]:
        out_dir = (assets_root / "converted" / file_id) if assets_root else (Path(p.parent) / ".converted")
        converted = _libreoffice_convert(p, out_dir, "pptx")
        if converted:
            docs = parse_file(
                str(converted),
                root_dir=str(converted.parent),
                source_type=source_type,
                session_id=session_id,
                assets_dir=assets_dir,
                parse_cfg=parse_cfg,
            )
            for d in docs:
                d.meta["orig_abs_path"] = d.meta.get("abs_path")
                d.meta["orig_rel_path"] = d.meta.get("rel_path")
                d.meta["abs_path"] = str(p.resolve())
                d.meta["rel_path"] = rel
                d.meta["source"] = rel
                # keep slide loc for citation
                if "slide" in d.meta:
                    d.meta["loc"] = f"{rel}#slide={d.meta['slide']}"
                else:
                    d.meta["loc"] = rel
                d.meta["file_ext"] = ext
            return docs
        return []

    if ext in [".xlsx", ".xlsm"]:
        try:
            import openpyxl  # type: ignore
        except Exception as e:
            raise RuntimeError("openpyxl is required: pip install openpyxl") from e

        excel_cfg = (parse_cfg.get("excel") or {})
        max_rows = int(excel_cfg.get("max_rows", 200))
        max_cols = int(excel_cfg.get("max_cols", 30))

        wb = openpyxl.load_workbook(str(p), data_only=True)
        docs: List[ParsedDocument] = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            lines: List[str] = []
            for r in ws.iter_rows(min_row=1, max_row=max_rows, max_col=max_cols, values_only=True):
                row = [str(x).strip() for x in (r or []) if x is not None and str(x).strip()]
                if row:
                    lines.append("\t".join(row))
            txt = "\n".join(lines).strip()
            if not txt:
                continue
            meta = dict(base_meta)
            meta.update({
                "sheet": sheet_name,
                "part": f"sheet_{sheet_name}",
                "loc": f"{rel}#sheet={sheet_name}",
            })
            docs.append(ParsedDocument(text=txt, meta=meta))
        return docs

    if ext in [".xls"]:
        out_dir = (assets_root / "converted" / file_id) if assets_root else (Path(p.parent) / ".converted")
        converted = _libreoffice_convert(p, out_dir, "xlsx")
        if converted:
            docs = parse_file(
                str(converted),
                root_dir=str(converted.parent),
                source_type=source_type,
                session_id=session_id,
                assets_dir=assets_dir,
                parse_cfg=parse_cfg,
            )
            for d in docs:
                d.meta["orig_abs_path"] = d.meta.get("abs_path")
                d.meta["orig_rel_path"] = d.meta.get("rel_path")
                d.meta["abs_path"] = str(p.resolve())
                d.meta["rel_path"] = rel
                d.meta["source"] = rel
                if "sheet" in d.meta:
                    d.meta["loc"] = f"{rel}#sheet={d.meta['sheet']}"
                else:
                    d.meta["loc"] = rel
                d.meta["file_ext"] = ext
            return docs
        return []

    if ext in [".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tiff"]:
        img_cfg = (parse_cfg.get("image") or {})
        do_ocr = bool(img_cfg.get("ocr", True))
        use_vlm_if_sparse = bool(img_cfg.get("use_vlm_if_ocr_sparse", True))
        sparse_chars = int(img_cfg.get("ocr_sparse_chars", 30))
        detect_chart_like = bool(img_cfg.get("detect_chart_like", True))
        chart_line_thr = int(img_cfg.get("chart_like_hough_lines", 25))
        vlm_prompt = str(img_cfg.get("vlm_prompt", "请描述这张图片中的关键信息，用于教学检索。"))

        vlm_cfg = (parse_cfg.get("vlm") or {})

        txt = ""
        meta = dict(base_meta)
        meta.update({"loc": rel, "part": "image", "image_path": str(p.resolve())})

        if do_ocr:
            try:
                o = ocr_image(str(p))
                txt = (o.text or "").strip()
                if assets_root:
                    out_dir = assets_root / "ocr" / file_id
                    out_dir.mkdir(parents=True, exist_ok=True)
                    ocr_path = out_dir / "image.ocr.json"
                    save_ocr_result(o, str(ocr_path))
                    meta.update({"ocr_path": str(ocr_path), "ocr_engine": o.engine, "ocr_mean_conf": o.mean_conf})
            except Exception:
                pass

        # Optional VLM route (OCR sparse or chart-like)
        need_vlm = False
        if use_vlm_if_sparse and (len(txt) < sparse_chars):
            need_vlm = True
        if (not need_vlm) and detect_chart_like:
            try:
                if _is_chart_like(str(p), line_threshold=chart_line_thr):
                    need_vlm = True
                    meta["chart_like"] = True
            except Exception:
                pass

        if need_vlm:
            try:
                vr = describe_image(str(p), vlm_prompt, vlm_cfg)
            except Exception:
                vr = None
            if vr is not None and (vr.text or "").strip():
                vtxt = (vr.text or "").strip()
                if assets_root:
                    out_dir = assets_root / "vlm" / file_id
                    out_dir.mkdir(parents=True, exist_ok=True)
                    vlm_path = out_dir / "image.vlm.json"
                    save_vlm_result(vr, str(vlm_path))
                    meta.update({"vlm_path": str(vlm_path), "vlm_engine": vr.engine})
                meta["vlm_used"] = True
                if txt:
                    txt = f"【OCR文本】\n{txt}\n\n【图像理解】\n{vtxt}".strip()
                else:
                    txt = f"【图像理解】\n{vtxt}".strip()

        if not txt:
            return []
        return [ParsedDocument(text=txt, meta=meta)]

    # ---------- audio / video ----------
    if ext in [".wav", ".mp3", ".m4a", ".flac", ".aac", ".ogg", ".opus"]:
        audio_cfg = (parse_cfg.get("audio") or {})
        prefer = str(audio_cfg.get("prefer", "faster-whisper"))
        model_size = str(audio_cfg.get("model_size", "small"))
        segment_level = bool(audio_cfg.get("segment_level", True))

        if not has_ffmpeg():
            raise RuntimeError("ffmpeg is required for audio/video parsing.")

        # ensure wav
        if ext != ".wav":
            if assets_root:
                out_dir = assets_root / "audio" / file_id
                out_dir.mkdir(parents=True, exist_ok=True)
                wav_path = out_dir / "audio.wav"
            else:
                wav_path = p.with_suffix(".wav")
            if not wav_path.exists():
                convert_audio_to_wav(str(p), str(wav_path))
            audio_path = wav_path
        else:
            audio_path = p

        asr = transcribe_audio(str(audio_path), prefer=prefer, model_size=model_size)
        docs: List[ParsedDocument] = []

        if assets_root:
            out_dir = assets_root / "asr" / file_id
            out_dir.mkdir(parents=True, exist_ok=True)
            asr_path = out_dir / "asr.json"
            save_asr_result(asr, str(asr_path))
        else:
            asr_path = None

        if segment_level and asr.segments:
            for j, seg in enumerate(asr.segments, start=1):
                t = (seg.text or "").strip()
                if not t:
                    continue
                meta = dict(base_meta)
                meta.update({
                    "part": f"asr_seg_{j}",
                    "loc": f"{rel}#t={seg.start:.2f}-{seg.end:.2f}",
                    "audio_path": str(p.resolve()),
                    "asr_engine": asr.engine,
                    "asr_path": str(asr_path) if asr_path else None,
                    "frame_ts": int(seg.start),
                    "start": float(seg.start),
                    "end": float(seg.end),
                })
                docs.append(ParsedDocument(text=t, meta=meta))
        else:
            meta = dict(base_meta)
            meta.update({
                "part": "asr",
                "loc": rel,
                "audio_path": str(p.resolve()),
                "asr_engine": asr.engine,
                "asr_path": str(asr_path) if asr_path else None,
            })
            if asr.text.strip():
                docs.append(ParsedDocument(text=asr.text.strip(), meta=meta))

        return docs

    if ext in [".mp4", ".mov", ".mkv", ".avi", ".webm"]:
        video_cfg = (parse_cfg.get("video") or {})
        prefer = str(video_cfg.get("prefer", "faster-whisper"))
        model_size = str(video_cfg.get("model_size", "small"))
        segment_level = bool(video_cfg.get("segment_level", True))

        extract_frames = bool(video_cfg.get("extract_frames", True))
        frame_interval = float(video_cfg.get("frame_interval_sec", 5.0))
        max_frames = int(video_cfg.get("max_frames", 20))
        ocr_frames = bool(video_cfg.get("ocr_frames", False))

        # VLM frame describing (optional)
        describe_frames = bool(video_cfg.get("describe_frames_with_vlm", True))
        vlm_max_frames = int(video_cfg.get("vlm_max_frames", 8))
        vlm_if_ocr_sparse = bool(video_cfg.get("vlm_if_ocr_sparse", True))
        sparse_chars = int(video_cfg.get("ocr_sparse_chars", 30))
        vlm_prompt = str(video_cfg.get("vlm_prompt", "请详细描述这张图片中的教学内容、图表信息和关键文字。"))
        combine_asr_and_frames = bool(video_cfg.get("combine_asr_and_frames", True))

        vlm_cfg = (parse_cfg.get("vlm") or {})

        if not has_ffmpeg():
            raise RuntimeError("ffmpeg is required for audio/video parsing.")

        # assets dirs
        if assets_root:
            vdir = assets_root / "video" / file_id
            vdir.mkdir(parents=True, exist_ok=True)
            wav_path = vdir / "audio.wav"
            frame_dir = vdir / "frames"
        else:
            vdir = None
            wav_path = p.with_suffix(".wav")
            frame_dir = p.parent / ".frames"

        if not wav_path.exists():
            extract_audio_with_ffmpeg(str(p), str(wav_path))

        asr = transcribe_audio(str(wav_path), prefer=prefer, model_size=model_size)

        asr_path = None
        if assets_root and vdir is not None:
            asr_path = vdir / "asr.json"
            save_asr_result(asr, str(asr_path))

        frame_paths: List[str] = []
        if extract_frames:
            try:
                frame_paths = extract_frames_with_ffmpeg(str(p), str(frame_dir), interval_sec=frame_interval, max_frames=max_frames)
            except Exception:
                frame_paths = []

        docs: List[ParsedDocument] = []

        # Build a helper to pick nearest frame by start time.
        def _pick_frame(ts: float) -> Optional[str]:
            if not frame_paths:
                return None
            idx = int(max(0, round(ts / frame_interval)))
            idx = min(idx, len(frame_paths) - 1)
            return frame_paths[idx]

        if segment_level and asr.segments:
            for j, seg in enumerate(asr.segments, start=1):
                t = (seg.text or "").strip()
                if not t:
                    continue
                meta = dict(base_meta)
                meta.update({
                    "part": f"asr_seg_{j}",
                    "loc": f"{rel}#t={seg.start:.2f}-{seg.end:.2f}",
                    "video_path": str(p.resolve()),
                    "asr_engine": asr.engine,
                    "asr_path": str(asr_path) if asr_path else None,
                    "frame_ts": int(seg.start),
                    "start": float(seg.start),
                    "end": float(seg.end),
                })
                fp = _pick_frame(seg.start)
                if fp:
                    meta["frame_path"] = fp
                docs.append(ParsedDocument(text=t, meta=meta))
        else:
            meta = dict(base_meta)
            meta.update({
                "part": "asr",
                "loc": rel,
                "video_path": str(p.resolve()),
                "asr_engine": asr.engine,
                "asr_path": str(asr_path) if asr_path else None,
            })
            if asr.text.strip():
                docs.append(ParsedDocument(text=asr.text.strip(), meta=meta))

        # Optional OCR on frames as extra evidence docs
        if ocr_frames and frame_paths:
            for i, fp in enumerate(frame_paths, start=1):
                try:
                    o = ocr_image(fp)
                    txt = (o.text or "").strip()
                    if not txt:
                        continue
                    meta = dict(base_meta)
                    meta.update({
                        "part": f"frame_{i}_ocr",
                        "loc": f"{rel}#frame={i}",
                        "frame_path": fp,
                        "frame_ts": int((i - 1) * frame_interval),
                        "ocr_engine": o.engine,
                        "ocr_mean_conf": o.mean_conf,
                    })
                    if assets_root:
                        ocr_dir = Path(fp).parent
                        ocr_path = Path(fp).with_suffix(".ocr.json")
                        save_ocr_result(o, str(ocr_path))
                        meta["ocr_path"] = str(ocr_path)
                    docs.append(ParsedDocument(text=txt, meta=meta))
                except Exception:
                    continue

        # Optional VLM description on key frames (aligned with模块二做法：OCR→必要时VLM)
        frame_desc_items: List[tuple[int, str, Dict[str, Any]]] = []  # (frame_idx, text, meta)
        if describe_frames and frame_paths:
            # choose a subset (evenly spaced)
            idxs: List[int]
            if len(frame_paths) <= vlm_max_frames:
                idxs = list(range(len(frame_paths)))
            else:
                step = max(1, int(len(frame_paths) / vlm_max_frames))
                idxs = list(range(0, len(frame_paths), step))[:vlm_max_frames]

            for i0 in idxs:
                fp = frame_paths[i0]
                frame_idx = i0 + 1
                ts = int(i0 * frame_interval)

                # route: run OCR first if requested
                ocr_txt = ""
                if vlm_if_ocr_sparse:
                    try:
                        o = ocr_image(fp)
                        ocr_txt = (o.text or "").strip()
                    except Exception:
                        ocr_txt = ""

                need_vlm = True
                if vlm_if_ocr_sparse:
                    need_vlm = len(ocr_txt) < sparse_chars

                if not need_vlm:
                    # OCR already rich enough; optionally store as a frame note
                    continue

                try:
                    vr = describe_image(fp, vlm_prompt, vlm_cfg)
                except Exception:
                    vr = None
                if vr is None or not (vr.text or "").strip():
                    continue

                vtxt = (vr.text or "").strip()
                meta = dict(base_meta)
                meta.update({
                    "part": f"frame_{frame_idx}_vlm",
                    "loc": f"{rel}#frame={frame_idx}",
                    "frame_path": fp,
                    "frame_ts": ts,
                    "vlm_engine": vr.engine,
                    "vlm_used": True,
                })
                if assets_root and vdir is not None:
                    out_dir = assets_root / "vlm" / file_id
                    out_dir.mkdir(parents=True, exist_ok=True)
                    vlm_path = out_dir / f"frame_{frame_idx:04d}.vlm.json"
                    save_vlm_result(vr, str(vlm_path))
                    meta["vlm_path"] = str(vlm_path)

                docs.append(ParsedDocument(text=f"【关键帧@{ts}s】\n{vtxt}".strip(), meta=meta))
                frame_desc_items.append((frame_idx, vtxt, meta))

        # Optional combined video evidence: ASR + frame descriptions
        if combine_asr_and_frames and (asr.text or "").strip() and frame_desc_items:
            lines: List[str] = []
            lines.append("【视频解说词（ASR）】")
            lines.append((asr.text or "").strip())
            lines.append("")
            lines.append("【关键帧描述（VLM）】")
            for frame_idx, vtxt, _m in frame_desc_items:
                ts = int((frame_idx - 1) * frame_interval)
                lines.append(f"- t={ts}s (frame={frame_idx}): {vtxt}")
            combined = "\n".join(lines).strip()

            meta = dict(base_meta)
            meta.update({
                "part": "video_combined",
                "loc": f"{rel}#combined",
                "video_path": str(p.resolve()),
                "asr_engine": asr.engine,
                "asr_path": str(asr_path) if asr_path else None,
                "frames_used": len(frame_desc_items),
            })
            docs.append(ParsedDocument(text=combined, meta=meta))

        return docs

    # Unknown types: try plain text
    txt = _safe_read_text(p).strip()
    if txt:
        meta = dict(base_meta)
        meta.update({"loc": rel, "part": "text_fallback"})
        return [ParsedDocument(text=txt, meta=meta)]
    return []

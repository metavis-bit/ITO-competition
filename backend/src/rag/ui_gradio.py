"""
多模态 AI 互动式教学智能体 — 最终版 UI（Premium 视觉重构 + 光标修复）
"""
from __future__ import annotations
import json, traceback
from pathlib import Path
from typing import List, Optional
import gradio as gr
from .service import RAGService

# ═══════════════════════════════════════════
# 强制 Light 主题（解决深色背景 + 光标不可见）
# ═══════════════════════════════════════════
try:
    THEME = gr.themes.Soft(
        primary_hue=gr.themes.colors.indigo,
        secondary_hue=gr.themes.colors.emerald,
        neutral_hue=gr.themes.colors.slate,
        font=gr.themes.GoogleFont("Outfit"),
        font_mono=gr.themes.GoogleFont("JetBrains Mono"),
    ).set(
        body_background_fill="#f8f6f3",
        body_background_fill_dark="#f8f6f3",
        block_background_fill="white",
        block_background_fill_dark="white",
        block_border_width="1px",
        block_border_color="rgba(0,0,0,0.06)",
        block_radius="24px",
        input_background_fill="#faf8f5",
        input_background_fill_dark="#faf8f5",
        input_border_color="rgba(0,0,0,0.06)",
        input_border_color_dark="rgba(0,0,0,0.06)",
        input_border_color_focus="#818cf8",
        input_border_color_focus_dark="#818cf8",
        input_border_width="1.5px",
        input_radius="10px",
        input_shadow="none",
        input_shadow_focus="0 0 0 3px rgba(99,102,241,0.12), 0 4px 20px rgba(99,102,241,0.12)",
        button_primary_background_fill="linear-gradient(135deg, #6366f1, #4f46e5)",
        button_primary_background_fill_dark="linear-gradient(135deg, #6366f1, #4f46e5)",
        button_primary_background_fill_hover="linear-gradient(135deg, #818cf8, #6366f1)",
        button_primary_text_color="white",
        button_primary_border_color="transparent",
        button_primary_shadow="0 4px 14px rgba(99,102,241,0.3)",
        button_secondary_background_fill="rgba(255,255,255,0.72)",
        button_secondary_background_fill_dark="rgba(255,255,255,0.72)",
        button_secondary_text_color="#64748b",
        button_secondary_border_color="rgba(0,0,0,0.06)",
        checkbox_background_color="white",
        checkbox_background_color_dark="white",
        slider_color="#6366f1",
        chatbot_code_background_color="#f1f5f9",
    )
except Exception:
    THEME = None

# ═══════════════════════════════════════════
# Premium CSS — 暖色编辑风 + 玻璃拟态 + 微动效 + 光标修复
# ═══════════════════════════════════════════
CUSTOM_CSS = """
@import url('https://fonts.googleapis.com/css2?family=Outfit:wght@300;400;500;600;700;800;900&family=Noto+Sans+SC:wght@300;400;500;700;900&display=swap');

/* ══ 强制浅色模式（修复深色主题 + 光标不可见）══ */
:root, .dark {
  --body-background-fill: #f8f6f3 !important;
  --block-background-fill: #ffffff !important;
  --input-background-fill: #faf8f5 !important;
  --body-text-color: #1a1a2e !important;
  --block-label-text-color: #475569 !important;
  --block-title-text-color: #1a1a2e !important;
  --neutral-50: #f8fafc !important;
  --neutral-100: #f1f5f9 !important;
  --neutral-200: #e2e8f0 !important;
  --neutral-300: #cbd5e1 !important;
  --neutral-400: #94a3b8 !important;
  --neutral-500: #64748b !important;
  --neutral-600: #475569 !important;
  --neutral-700: #334155 !important;
  --neutral-800: #1e293b !important;
  --neutral-900: #0f172a !important;
  --neutral-950: #020617 !important;
  --background-fill-primary: #ffffff !important;
  --background-fill-secondary: #f8f6f3 !important;
  --border-color-primary: rgba(0,0,0,0.06) !important;
  --border-color-accent: rgba(99,102,241,0.15) !important;
  --chatbot-body-background-fill: #ffffff !important;
  --chatbot-bubble-bot-background-fill: #f1f5f9 !important;
  --chatbot-bubble-user-background-fill: rgba(99,102,241,0.08) !important;
  --chatbot-bubble-bot-text-color: #1a1a2e !important;
  --chatbot-bubble-user-text-color: #1a1a2e !important;
  --panel-background-fill: #ffffff !important;
  --table-even-background-fill: #faf8f5 !important;
  --table-odd-background-fill: #ffffff !important;
  --color-accent-soft: rgba(99,102,241,0.12) !important;
  color-scheme: light !important;
}

/* ── 全局色板 ── */
:root {
  --bg-base: #f8f6f3;
  --bg-warm: #faf8f5;
  --surface: rgba(255, 255, 255, 0.72);
  --surface-hover: rgba(255, 255, 255, 0.88);
  --border-soft: rgba(0, 0, 0, 0.06);
  --border-accent: rgba(99, 102, 241, 0.15);
  --text-primary: #1a1a2e;
  --text-secondary: #64748b;
  --text-muted: #94a3b8;
  --accent: #6366f1;
  --accent-light: #818cf8;
  --accent-glow: rgba(99, 102, 241, 0.12);
  --accent-deep: #4f46e5;
  --success: #10b981;
  --success-glow: rgba(16, 185, 129, 0.1);
  --warning: #f59e0b;
  --danger: #ef4444;
  --shadow-sm: 0 1px 3px rgba(0,0,0,0.04), 0 1px 2px rgba(0,0,0,0.03);
  --shadow-md: 0 4px 20px rgba(0,0,0,0.06), 0 2px 8px rgba(0,0,0,0.04);
  --shadow-lg: 0 12px 40px rgba(0,0,0,0.08), 0 4px 12px rgba(0,0,0,0.04);
  --shadow-glow: 0 0 0 3px var(--accent-glow), 0 4px 20px rgba(99,102,241,0.12);
  --radius-sm: 10px;
  --radius-md: 16px;
  --radius-lg: 24px;
  --radius-xl: 32px;
}

/* ── 全局基底 ── */
.gradio-container, .gradio-container.dark {
  font-family: 'Outfit', 'Noto Sans SC', system-ui, sans-serif !important;
  background: var(--bg-base) !important;
  background-image:
    radial-gradient(ellipse 80% 60% at 10% 0%, rgba(99,102,241,0.06) 0%, transparent 60%),
    radial-gradient(ellipse 60% 50% at 90% 100%, rgba(16,185,129,0.05) 0%, transparent 60%) !important;
  min-height: 100vh;
  color: var(--text-primary) !important;
}

/* ── 隐藏默认footer ── */
footer { display: none !important; }

/* ══ 强制所有输入框浅色 + 光标可见 ══ */
.gradio-container textarea,
.gradio-container input[type="text"],
.gradio-container input[type="search"],
.gradio-container input[type="number"],
.gradio-container .wrap input,
.gradio-container .wrap textarea,
.dark textarea,
.dark input[type="text"],
.dark input[type="search"],
.dark input[type="number"] {
  background: var(--bg-warm) !important;
  color: var(--text-primary) !important;
  caret-color: var(--accent) !important;
  font-family: 'Outfit', 'Noto Sans SC', sans-serif !important;
  border: 1.5px solid var(--border-soft) !important;
  border-radius: var(--radius-sm) !important;
  transition: all 0.2s ease !important;
}
.gradio-container textarea:focus,
.gradio-container input[type="text"]:focus,
.dark textarea:focus,
.dark input[type="text"]:focus {
  border-color: var(--accent-light) !important;
  box-shadow: var(--shadow-glow) !important;
  outline: none !important;
}
.gradio-container textarea::placeholder,
.gradio-container input::placeholder {
  color: var(--text-muted) !important;
  opacity: 1 !important;
}

/* ══ 下拉框 / Checkbox / Label 强制浅色 ══ */
.gradio-container select,
.gradio-container .wrap select,
.gradio-container [data-testid="dropdown"],
.gradio-container .secondary-wrap,
.dark select {
  background: #ffffff !important;
  color: var(--text-primary) !important;
}
.gradio-container label span,
.gradio-container .label-wrap span {
  color: var(--text-secondary) !important;
}

/* ── 顶部 Hero Banner ── */
.hero-banner {
  background: linear-gradient(135deg, #1a1a2e 0%, #2d1b69 40%, #44337a 100%) !important;
  border-radius: var(--radius-xl) !important;
  padding: 36px 40px !important;
  margin-bottom: 8px !important;
  position: relative;
  overflow: hidden;
  border: none !important;
}
.hero-banner::before {
  content: '';
  position: absolute;
  top: -50%; right: -20%;
  width: 60%; height: 200%;
  background: radial-gradient(circle, rgba(99,102,241,0.15) 0%, transparent 70%);
  animation: hero-pulse 8s ease-in-out infinite alternate;
  pointer-events: none;
}
.hero-banner::after {
  content: '';
  position: absolute;
  bottom: -30%; left: -10%;
  width: 40%; height: 160%;
  background: radial-gradient(circle, rgba(16,185,129,0.1) 0%, transparent 70%);
  animation: hero-pulse 6s ease-in-out infinite alternate-reverse;
  pointer-events: none;
}
@keyframes hero-pulse {
  0% { transform: scale(1) translate(0, 0); }
  100% { transform: scale(1.1) translate(2%, -2%); }
}

/* ── 玻璃面板 ── */
.glass-panel {
  background: var(--surface) !important;
  backdrop-filter: blur(20px) saturate(1.2) !important;
  -webkit-backdrop-filter: blur(20px) saturate(1.2) !important;
  border: 1px solid var(--border-soft) !important;
  border-radius: var(--radius-lg) !important;
  box-shadow: var(--shadow-md) !important;
  padding: 24px !important;
  transition: box-shadow 0.3s ease, border-color 0.3s ease !important;
  animation: fadeInUp 0.5s ease-out both;
}
.glass-panel:hover {
  box-shadow: var(--shadow-lg) !important;
  border-color: var(--border-accent) !important;
}

/* ── 按钮系统（3D 按下弹跳效果）── */
button.primary, .primary-btn {
  background: linear-gradient(135deg, var(--accent) 0%, var(--accent-deep) 100%) !important;
  color: white !important;
  font-family: 'Outfit', 'Noto Sans SC', sans-serif !important;
  font-weight: 800 !important;
  font-size: 1.05rem !important;
  border-radius: var(--radius-md) !important;
  border: none !important;
  padding: 12px 28px !important;
  box-shadow: 0 4px 0 #3730a3, 0 6px 16px rgba(99,102,241,0.25) !important;
  transition: all 0.15s cubic-bezier(0.4, 0, 0.2, 1) !important;
  text-transform: uppercase;
  letter-spacing: 1px;
  position: relative;
  overflow: hidden;
}
button.primary:hover, .primary-btn:hover {
  background: linear-gradient(135deg, var(--accent-light) 0%, var(--accent) 100%) !important;
  transform: translateY(-1px) !important;
  box-shadow: 0 5px 0 #3730a3, 0 8px 20px rgba(99,102,241,0.3) !important;
}
button.primary:active, .primary-btn:active {
  transform: translateY(4px) !important;
  box-shadow: 0 0 0 #3730a3, 0 2px 8px rgba(99,102,241,0.2) !important;
}

button.secondary, .secondary-btn {
  background: #ffffff !important;
  color: var(--text-secondary) !important;
  font-family: 'Outfit', 'Noto Sans SC', sans-serif !important;
  font-weight: 700 !important;
  border-radius: var(--radius-md) !important;
  border: 2px solid #e2e8f0 !important;
  box-shadow: 0 4px 0 #e2e8f0 !important;
  transition: all 0.15s cubic-bezier(0.4, 0, 0.2, 1) !important;
}
button.secondary:hover, .secondary-btn:hover {
  background: #faf8f5 !important;
  border-color: var(--accent-light) !important;
  color: var(--accent) !important;
  box-shadow: 0 4px 0 var(--accent-light) !important;
  transform: translateY(-1px) !important;
}
button.secondary:active, .secondary-btn:active {
  transform: translateY(4px) !important;
  box-shadow: 0 0 0 #e2e8f0 !important;
}

/* ── 生成按钮特殊高亮（绿色 3D 弹跳）── */
.generate-btn button {
  background: linear-gradient(135deg, var(--success) 0%, #059669 100%) !important;
  box-shadow: 0 4px 0 #047857, 0 6px 16px rgba(16,185,129,0.25) !important;
  font-size: 1.05rem !important;
  font-weight: 800 !important;
  padding: 14px 32px !important;
  text-transform: uppercase !important;
  letter-spacing: 0.8px !important;
  transition: all 0.15s cubic-bezier(0.4, 0, 0.2, 1) !important;
}
.generate-btn button:hover {
  transform: translateY(-1px) !important;
  box-shadow: 0 5px 0 #047857, 0 8px 20px rgba(16,185,129,0.3) !important;
  background: linear-gradient(135deg, #34d399 0%, #10b981 100%) !important;
}
.generate-btn button:active {
  transform: translateY(4px) !important;
  box-shadow: 0 0 0 #047857, 0 2px 8px rgba(16,185,129,0.2) !important;
}

/* ── Chatbot 样式（全面覆盖深色）── */
.chatbot-container,
.chatbot-container > div,
.chatbot-container .wrap,
.chatbot-container .bot,
.chatbot-container .user,
[data-testid="chatbot"],
[data-testid="chatbot"] > div {
  background: #ffffff !important;
  border-radius: var(--radius-lg) !important;
  border: 1px solid var(--border-soft) !important;
  overflow: hidden !important;
}
/* 气泡文字强制深色 */
.chatbot-container .message,
.chatbot-container .bot .message-wrap,
.chatbot-container .user .message-wrap,
.chatbot-container .message-wrap,
[data-testid="chatbot"] .message,
[data-testid="chatbot"] .bot,
[data-testid="chatbot"] .user,
[data-testid="chatbot"] .message-wrap {
  font-family: 'Outfit', 'Noto Sans SC', sans-serif !important;
  border-radius: var(--radius-md) !important;
  color: #1a1a2e !important;
}
/* Bot 气泡浅灰底 */
.chatbot-container .bot .message-wrap,
[data-testid="chatbot"] .bot .message-wrap,
[data-testid="chatbot"] .bot > div {
  background: #f1f5f9 !important;
  color: #1a1a2e !important;
}
/* User 气泡淡紫底 */
.chatbot-container .user .message-wrap,
[data-testid="chatbot"] .user .message-wrap,
[data-testid="chatbot"] .user > div {
  background: rgba(99,102,241,0.08) !important;
  color: #1a1a2e !important;
}
/* 气泡内所有文本 */
.chatbot-container p, .chatbot-container span, .chatbot-container li,
.chatbot-container code, .chatbot-container pre, .chatbot-container h1,
.chatbot-container h2, .chatbot-container h3, .chatbot-container h4,
[data-testid="chatbot"] p, [data-testid="chatbot"] span,
[data-testid="chatbot"] li, [data-testid="chatbot"] code {
  color: #1a1a2e !important;
}

/* ── Audio 组件强制浅色 ── */
.gradio-container audio,
.gradio-container .audio-container,
.gradio-container [data-testid="audio"],
.gradio-container .wrap.svelte-aqlk7e,
.gradio-container .audio-player,
.dark audio, .dark .audio-container {
  background: #f0eef8 !important;
  color: #1a1a2e !important;
  border-radius: var(--radius-md) !important;
  border: 1.5px solid var(--border-soft) !important;
}

/* ══ 全局 div/span/p 兜底（防止任何漏网深色）══ */
.gradio-container .block,
.gradio-container .form,
.gradio-container .wrap,
.gradio-container .panel,
.dark .block,
.dark .form,
.dark .wrap {
  background-color: transparent !important;
  color: var(--text-primary) !important;
}
.gradio-container .block.padded,
.dark .block.padded {
  background: #ffffff !important;
}

/* ── Tab 导航 ── */
.tabs > div > button {
  font-family: 'Outfit', 'Noto Sans SC', sans-serif !important;
  font-weight: 600 !important;
  font-size: 0.95rem !important;
  border-radius: var(--radius-sm) var(--radius-sm) 0 0 !important;
  transition: all 0.2s ease !important;
  padding: 10px 20px !important;
  color: var(--text-secondary) !important;
  background: transparent !important;
}
.tabs > div > button.selected {
  color: var(--accent) !important;
  border-bottom: 3px solid var(--accent) !important;
  background: var(--accent-glow) !important;
}

/* ── 预览面板 ── */
.preview-html {
  border-radius: var(--radius-lg) !important;
  overflow: hidden !important;
  background: white !important;
  box-shadow: var(--shadow-md) !important;
  border: 1px solid var(--border-soft) !important;
  min-height: 480px;
  transition: box-shadow 0.3s ease !important;
}
.preview-html:hover {
  box-shadow: var(--shadow-lg) !important;
}

/* ── Accordion ── */
.gradio-accordion {
  border-radius: var(--radius-md) !important;
  border: 1px solid var(--border-soft) !important;
  overflow: hidden !important;
}

/* ── 状态指示文本 ── */
.status-box textarea {
  font-size: 0.88rem !important;
  color: var(--text-secondary) !important;
  background: linear-gradient(135deg, var(--accent-glow), var(--success-glow)) !important;
  border: 1px dashed var(--border-accent) !important;
  border-radius: var(--radius-sm) !important;
  caret-color: var(--accent) !important;
}

/* ── Section 标题 ── */
.section-title {
  font-size: 1.1rem !important;
  font-weight: 700 !important;
  color: var(--text-primary) !important;
  margin-bottom: 4px !important;
  letter-spacing: -0.3px;
}

/* ── 文件上传区 ── */
.gradio-file, [data-testid="dropzone"] {
  border-radius: var(--radius-md) !important;
  border: 2px dashed var(--border-soft) !important;
  background: var(--bg-warm) !important;
  transition: border-color 0.2s ease !important;
}
.gradio-file:hover, [data-testid="dropzone"]:hover {
  border-color: var(--accent-light) !important;
}

/* ── Slider ── */
.gradio-container input[type="range"] {
  accent-color: var(--accent) !important;
}

/* ── 滚动条美化 ── */
::-webkit-scrollbar { width: 6px; height: 6px; }
::-webkit-scrollbar-track { background: transparent; }
::-webkit-scrollbar-thumb { background: rgba(0,0,0,0.12); border-radius: 3px; }
::-webkit-scrollbar-thumb:hover { background: rgba(0,0,0,0.2); }

/* ── 淡入动画 ── */
@keyframes fadeInUp {
  from { opacity: 0; transform: translateY(12px); }
  to { opacity: 1; transform: translateY(0); }
}
"""

_empty = """
<div style='display:flex;flex-direction:column;align-items:center;justify-content:center;height:100%;min-height:480px;
  color:#94a3b8;font-family:Outfit,Noto Sans SC,sans-serif;'>
  <div style="width:80px;height:80px;border-radius:20px;
    background:linear-gradient(135deg,rgba(99,102,241,0.08),rgba(16,185,129,0.08));
    display:flex;align-items:center;justify-content:center;margin-bottom:20px;">
    <svg width="36" height="36" viewBox="0 0 24 24" fill="none" stroke="#6366f1" stroke-width="1.5" stroke-linecap="round">
      <rect x="2" y="3" width="20" height="14" rx="2.5"></rect>
      <line x1="8" y1="21" x2="16" y2="21"></line><line x1="12" y1="17" x2="12" y2="21"></line>
      <circle cx="12" cy="10" r="2" fill="none" stroke="#10b981" stroke-width="1.5"/>
    </svg>
  </div>
  <h3 style="margin:0;font-size:1.3rem;font-weight:700;color:#334155;letter-spacing:-0.3px;">引擎就绪</h3>
  <p style="margin-top:10px;font-size:0.92rem;max-width:280px;text-align:center;line-height:1.7;color:#94a3b8;">
    在左侧输入教学需求<br/>AI 将自动生成完整课件
  </p>
</div>
"""

_HERO_MD = """
<div style="position:relative;z-index:1;display:flex;align-items:center;gap:20px;">
  <div style="flex-shrink:0;width:56px;height:56px;border-radius:16px;
    background:linear-gradient(135deg,rgba(255,255,255,0.2),rgba(255,255,255,0.05));
    backdrop-filter:blur(8px);border:1px solid rgba(255,255,255,0.15);
    display:flex;align-items:center;justify-content:center;">
    <svg width="28" height="28" viewBox="0 0 24 24" fill="none" stroke="white" stroke-width="2" stroke-linecap="round">
      <path d="M12 2L2 7l10 5 10-5-10-5z"/><path d="M2 17l10 5 10-5"/><path d="M2 12l10 5 10-5"/>
    </svg>
  </div>
  <div>
    <h1 style="margin:0;font-size:1.75rem;font-weight:800;color:white;letter-spacing:-0.5px;
      font-family:Outfit,Noto Sans SC,sans-serif;">
      多模态 AI 互动式教学智能体
    </h1>
    <p style="margin:6px 0 0;font-size:0.9rem;color:rgba(255,255,255,0.6);font-weight:400;
      font-family:Outfit,Noto Sans SC,sans-serif;letter-spacing:0.2px;">
      A04 锐捷网络 &nbsp;·&nbsp; 对话收集 → 智能生成 → 预览修改 → 迭代优化
    </p>
  </div>
</div>
"""


def build_ui(cfg_path: str = "config.yaml") -> gr.Blocks:
    service = RAGService(cfg_path)
    dialogue_mgr = None
    try:
        from .dialogue.dialogue_manager import DialogueManager
        dialogue_mgr = DialogueManager(cfg_path)
    except Exception:
        pass
    pipeline = None
    try:
        from .courseware_pipeline import CoursewarePipeline
        pipeline = CoursewarePipeline(cfg_path)
    except Exception:
        pass
    try:
        idxs = service.list_indexes() or ["kb"]
    except Exception:
        idxs = ["kb"]

    # ── 工具函数 ──
    def asr_to_text(audio_path) -> str:
        if not audio_path:
            return ""
        try:
            from .tools.asr import transcribe_audio
            r = transcribe_audio(str(audio_path))
            return str(r) if not isinstance(r, str) else r
        except Exception as e:
            return f"[语音识别需安装faster-whisper: {e}]"

    def pptx_preview(path):
        if not path or not Path(path).exists():
            return _empty
        try:
            from .ppt_preview import pptx_preview_html
            return pptx_preview_html(path, max_pages=12)
        except Exception:
            try:
                from pptx import Presentation
                prs = Presentation(path)
                lines = [
                    f"<p><b>第{i+1}页:</b> {' | '.join(s.text.strip()[:50] for s in sl.shapes if hasattr(s,'text') and s.text.strip())[:3]}</p>"
                    for i, sl in enumerate(prs.slides)
                ]
                return "<div>" + "\n".join(lines) + "</div>"
            except Exception as e:
                return f"<p>预览失败: {e}</p>"

    def docx_preview(path):
        if not path or not Path(path).exists():
            return "无Word文件"
        try:
            import docx
            doc = docx.Document(path)
            lines = []
            for p in doc.paragraphs[:30]:
                t = p.text.strip()
                if not t:
                    continue
                sn = p.style.name if p.style else ""
                if "Heading" in sn:
                    lines.append(f"\n## {t}")
                else:
                    lines.append(t)
            return "\n".join(lines)[:3000]
        except Exception as e:
            return f"预览失败: {e}"

    # ── 对话 ──
    def chat(user_msg, audio, history, files, note):
        if audio and not user_msg.strip():
            user_msg = asr_to_text(audio)
        if not user_msg.strip():
            return history or [], "", "请输入文字或语音", None
        if not dialogue_mgr:
            return (history or []) + [[user_msg, "⚠️ 对话模块未加载"]], "", "未就绪", None
        try:
            fps = [f.name if hasattr(f, 'name') else str(f) for f in (files if isinstance(files, list) else [files])] if files else []
            ns = [note] * len(fps) if fps else None
            reply, state = dialogue_mgr.chat(user_input=user_msg, uploaded_files=fps or None, file_notes=ns)
            ch = (history or []) + [[user_msg, reply]]
            info = (
                f"阶段: {state.value}\n"
                f"已收集:\n{dialogue_mgr.collected.to_text()}\n\n"
                f"缺失: {', '.join(dialogue_mgr.collected.missing_fields()) or '✅ 信息完整'}"
            )
            return ch, "", info, None
        except Exception as e:
            return (history or []) + [[user_msg, f"❌ {e}"]], "", str(e), None

    def reset():
        if dialogue_mgr:
            dialogue_mgr.reset()
        return [], "", "已重置", None

    # ── 一键生成 ──
    def gen_all(_st):
        if not pipeline:
            return "❌ 流水线未加载", _empty, "未生成", _empty, _empty, None, None, None, None, None
        if dialogue_mgr and not dialogue_mgr.collected.topic and not dialogue_mgr.history:
            return "⚠️ 请先在左侧对话中描述教学需求，然后再点此按钮。", _empty, "", _empty, _empty, None, None, None, None, None
        try:
            r = pipeline.generate_all(dialogue_mgr=dialogue_mgr, rag_service=service, output_types=["ppt", "docx", "game"])
            st = r.summary()
            ph = pptx_preview(r.pptx_path) if r.pptx_path else _empty
            dt = docx_preview(r.docx_path) if r.docx_path else "未生成"
            gh = r.game_html or _empty
            ah = r.animation_html or _empty
            pf = r.pptx_path if r.pptx_path and Path(r.pptx_path).exists() else None
            df = r.docx_path if r.docx_path and Path(r.docx_path).exists() else None
            gf = r.game_path if r.game_path and Path(r.game_path).exists() else None
            af = r.animation_path if r.animation_path and Path(r.animation_path).exists() else None
            return st, ph, dt, gh, ah, pf, df, gf, af, r
        except Exception as e:
            return f"❌ {e}\n{traceback.format_exc()}", _empty, "", _empty, _empty, None, None, None, None, None

    # ── 迭代修改 ──
    def regen(fb, rp, rd, rg, prev):
        if not pipeline:
            return "❌ 未加载", _empty, "", _empty, _empty, None, None, None, None, prev
        if not fb.strip():
            return "请输入修改意见", _empty, "", _empty, _empty, None, None, None, None, prev
        if not prev:
            return "请先生成课件", _empty, "", _empty, _empty, None, None, None, None, None
        rt = []
        if rp:
            rt.append("ppt")
        if rd:
            rt.append("docx")
        if rg:
            rt.append("game")
        if not rt:
            rt = ["ppt", "docx", "game"]
        try:
            r = pipeline.regenerate_with_feedback(feedback=fb, previous_result=prev, rag_service=service, regenerate_types=rt)
            st = "🔄 修改完成:\n" + r.summary()
            return (
                st,
                pptx_preview(r.pptx_path) if r.pptx_path else "",
                docx_preview(r.docx_path) if r.docx_path else "",
                r.game_html or "",
                r.animation_html or getattr(prev, 'animation_html', '') or "",
                r.pptx_path,
                r.docx_path,
                r.game_path,
                getattr(r, 'animation_path', None) or getattr(prev, 'animation_path', None),
                r,
            )
        except Exception as e:
            return f"❌ {e}", "", "", "", "", None, None, None, None, prev

    # ── RAG ──
    def query(q, idx, k, tr):
        try:
            res = service.query(question=q, indexes=idx or ["kb"], top_k=int(k), evidence_tag="ui", enable_trace=True)
            hv = res.get("human_view", {}) or {}
            return (
                hv.get("answer_md", ""),
                hv.get("evidence_md", ""),
                hv.get("sources_md", ""),
                hv.get("trace_md", "") if tr else "",
                json.dumps({"answer": res.get("answer")}, ensure_ascii=False, indent=2),
            )
        except Exception as e:
            return "", "", "", f"错误:{e}", ""

    # ═══════════════════════════════════════════
    # BUILD UI — theme + css 双保险
    # ═══════════════════════════════════════════
    blocks_kwargs = {"css": CUSTOM_CSS}
    if THEME is not None:
        blocks_kwargs["theme"] = THEME

    try:
        gr.Blocks(**blocks_kwargs)
        ok = True
    except Exception:
        ok = False

    if not ok:
        blocks_kwargs = {"css": CUSTOM_CSS}
        try:
            gr.Blocks(**blocks_kwargs)
        except Exception:
            blocks_kwargs = {}

    with gr.Blocks(**blocks_kwargs) as demo:
        cw = gr.State(None)

        # ── Hero Banner ──
        gr.Markdown(_HERO_MD, elem_classes="hero-banner")

        # ══════ Tab 1: 教学对话 & 课件生成 ══════
        with gr.Tab("✨ 智能共创中心"):
            with gr.Row(equal_height=False):
                # 左侧面板
                with gr.Column(scale=4, elem_classes="glass-panel"):
                    gr.Markdown("### 🎙️ 描述教学需求")
                    chatbot = gr.Chatbot(height=280, show_label=False, elem_classes="chatbot-container")
                    with gr.Row():
                        mi = gr.Textbox(
                            show_label=False,
                            placeholder="描述需求，如：帮我生成关于DNA转录的互动课...",
                            scale=5,
                            lines=1,
                            max_lines=3,
                        )
                        sb = gr.Button("发送 ✈️", elem_classes="secondary", scale=1)
                    au = gr.Audio(sources=["microphone"], type="filepath", label="🎙️ 语音输入")

                    with gr.Accordion("📎 上传参考资料", open=False):
                        gr.Markdown("<span style='color:#64748b;font-size:0.9em;'>上传PDF/Word/PPT/音视频，AI从中提取知识。每次上传请说明用途。</span>")
                        rf = gr.File(label="拖拽文件至此", file_count="multiple")
                        rn = gr.Textbox(label="资料用途说明", placeholder="如：参照此PDF第3章的格式和知识点")

                    si = gr.Textbox(label="💡 已收集信息", interactive=False, lines=3, elem_classes="status-box")

                    with gr.Row():
                        rb = gr.Button("🔄 重置对话", elem_classes="secondary")
                        gb = gr.Button("🚀 一键生成全部课件", variant="primary", elem_classes="generate-btn")
                    gs = gr.Textbox(label="⏳ 生成状态", interactive=False, lines=3, elem_classes="status-box")

                # 右侧预览面板
                with gr.Column(scale=7):
                    with gr.Tabs():
                        with gr.Tab("🎮 互动游戏"):
                            gp = gr.HTML(value=_empty, elem_classes="preview-html")
                            gd = gr.File(label="⏬ 下载游戏 (.html)")
                        with gr.Tab("🎬 知识动画"):
                            ap = gr.HTML(value=_empty, elem_classes="preview-html")
                            ad = gr.File(label="⏬ 下载动画 (.html)")
                        with gr.Tab("📊 PPT课件"):
                            pp = gr.HTML(value=_empty, elem_classes="preview-html")
                            pd = gr.File(label="⏬ 下载PPT (.pptx)")
                        with gr.Tab("📝 Word教案"):
                            dp = gr.Textbox(label="教案预览", lines=18, interactive=False, value="等待生成")
                            dd = gr.File(label="⏬ 下载教案 (.docx)")

            sb.click(fn=chat, inputs=[mi, au, chatbot, rf, rn], outputs=[chatbot, mi, si, au])
            mi.submit(fn=chat, inputs=[mi, au, chatbot, rf, rn], outputs=[chatbot, mi, si, au])
            rb.click(fn=reset, outputs=[chatbot, mi, si, au])
            gb.click(fn=gen_all, inputs=[cw], outputs=[gs, pp, dp, gp, ap, pd, dd, gd, ad, cw])

        # ══════ Tab 2: 迭代修改 ══════
        with gr.Tab("✏️ 迭代微调"):
            with gr.Row(equal_height=False):
                with gr.Column(scale=4, elem_classes="glass-panel"):
                    gr.Markdown("### 🎯 输入修改意见")
                    fi = gr.Textbox(
                        label="修改要求",
                        placeholder="如：把第3页简化 / 游戏改排序题 / 教案增加讨论环节",
                        lines=5,
                    )
                    gr.Markdown(
                        "<span style='color:#64748b;font-size:0.85em;'>选择需要重新生成的内容类型：</span>"
                    )
                    with gr.Row():
                        fp = gr.Checkbox(value=True, label="PPT")
                        fd = gr.Checkbox(value=True, label="Word")
                        fg = gr.Checkbox(value=True, label="游戏")
                    fb_btn = gr.Button("✨ 应用修改并重新生成", variant="primary", elem_classes="primary")
                    fs = gr.Textbox(label="状态", lines=2, interactive=False, elem_classes="status-box")

                with gr.Column(scale=7):
                    with gr.Tabs():
                        with gr.Tab("🎮 游戏"):
                            rgp = gr.HTML(value=_empty, elem_classes="preview-html")
                            rgd = gr.File(label="下载游戏")
                        with gr.Tab("🎬 动画"):
                            rap = gr.HTML(value=_empty, elem_classes="preview-html")
                            rad = gr.File(label="下载动画")
                        with gr.Tab("📊 PPT"):
                            rpp = gr.HTML(value=_empty, elem_classes="preview-html")
                            rpd = gr.File(label="下载 PPT")
                        with gr.Tab("📝 教案"):
                            rdp = gr.Textbox(lines=15, interactive=False)
                            rdd = gr.File(label="下载教案")

            fb_btn.click(fn=regen, inputs=[fi, fp, fd, fg, cw], outputs=[fs, rpp, rdp, rgp, rap, rpd, rdd, rgd, rad, cw])

        # ══════ Tab 3: 知识问答 ══════
        with gr.Tab("🔍 知识库问答"):
            with gr.Column(elem_classes="glass-panel"):
                gr.Markdown("### 💬 向知识库提问")
                q = gr.Textbox(label="提问", placeholder="向知识库提问...", lines=2)
                with gr.Row():
                    qi = gr.Dropdown(
                        choices=idxs,
                        value=["kb"] if "kb" in idxs else idxs[:1],
                        multiselect=True,
                        label="索引",
                    )
                    qk = gr.Slider(1, 20, value=5, step=1, label="Top-K")
                    qt = gr.Checkbox(value=False, label="显示Trace")
                qb = gr.Button("查询", variant="primary")
                qa = gr.Markdown()
                with gr.Accordion("📄 证据来源", open=False):
                    qe = gr.Markdown()
                    qs_md = gr.Markdown()
                with gr.Accordion("⚙️ Trace日志", open=False):
                    qtr = gr.Markdown()
                with gr.Accordion("JSON", open=False):
                    qj = gr.Code(language="json")
            qb.click(fn=query, inputs=[q, qi, qk, qt], outputs=[qa, qe, qs_md, qtr, qj])

        # ══════ Tab 4: 知识库管理 ══════
        with gr.Tab("📚 知识库管理"):
            with gr.Column(elem_classes="glass-panel"):
                ks = gr.Textbox(label="索引列表", interactive=False)
                gr.Button("🔄 刷新").click(
                    fn=lambda: "\n".join(f"  - {x}" for x in service.list_indexes()),
                    outputs=[ks],
                )
                gr.Markdown("---")
                with gr.Row():
                    id_ = gr.Textbox(label="资料目录", value="knowledge_base")
                    ii = gr.Textbox(label="索引名", value="kb")
                    iis = gr.Textbox(label="Session", value="kb")
                    ir = gr.Checkbox(value=True, label="重建索引")
                ib = gr.Button("📥 开始入库", variant="primary", elem_classes="primary")
                io = gr.Textbox(label="结果", lines=3, interactive=False)

                def do_ingest(d, i, s, r):
                    try:
                        if r:
                            service.engine.reset_index(i)
                        return f"✅ {service.engine.ingest_dir(d, index=i, source_type='knowledge_base', session_id=s or 'kb')}"
                    except Exception as e:
                        return f"❌ {e}"

                ib.click(fn=do_ingest, inputs=[id_, ii, iis, ir], outputs=[io])

    return demo


def main():
    demo = build_ui("config.yaml")
    demo.queue()
    demo.launch(server_name="127.0.0.1", server_port=7860, share=False)


if __name__ == "__main__":
    main()

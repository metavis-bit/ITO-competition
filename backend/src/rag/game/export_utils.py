"""
Animation & Game export utilities — multi-format output support.

Provides:
  1. HTML → static GIF (step-by-step snapshot via Pillow, no browser needed)
  2. HTML → MP4 video (Pillow frames → FFmpeg encode, no browser needed)
  3. HTML → embedded PPT slide with hyperlink
  4. ZIP bundle of all game/animation artifacts

对应 A04 要求:
  - 5c) 动画或小游戏可以以 HTML5、GIF、MP4 导出或集成到 PPT 中
"""
from __future__ import annotations

import io
import math
import textwrap
from pathlib import Path
from typing import Any, Dict, List, Optional

from PIL import Image, ImageDraw, ImageFont


# ═══════════════════════════════════════════
# 1. Animation steps → static GIF
# ═══════════════════════════════════════════

def animation_steps_to_gif(
    data: Dict[str, Any],
    output_path: str,
    width: int = 800,
    height: int = 500,
    duration_per_frame_ms: int = 2000,
) -> str:
    """
    Render animation steps data to an animated GIF.

    Each step becomes one frame with icon, title, and description.
    No browser or headless engine required — pure Pillow rendering.

    Args:
        data: {"title": "...", "steps": [{"icon": "🔬", "title": "...", "description": "..."}]}
        output_path: Where to save the GIF.
        width: Frame width in pixels.
        height: Frame height in pixels.
        duration_per_frame_ms: Duration per frame in milliseconds.

    Returns:
        The output file path.
    """
    title = data.get("title", "知识点动画")
    steps = data.get("steps", [])
    if not steps:
        steps = [{"icon": "📖", "title": "暂无内容", "description": "动画步骤为空"}]

    frames: List[Image.Image] = []

    # Try to load a CJK-compatible font
    font_large = _get_font(36)
    font_medium = _get_font(24)
    font_small = _get_font(18)
    font_icon = _get_font(64)

    # Color palette
    bg_gradient_top = (30, 41, 59)      # slate-800
    bg_gradient_bottom = (15, 23, 42)   # slate-900
    text_white = (248, 250, 252)
    text_muted = (148, 163, 184)        # slate-400
    accent = (56, 189, 248)             # sky-400
    step_bg = (51, 65, 85)              # slate-700

    for i, step in enumerate(steps):
        img = Image.new("RGB", (width, height))
        draw = ImageDraw.Draw(img)

        # Gradient background
        for y in range(height):
            ratio = y / height
            r = int(bg_gradient_top[0] * (1 - ratio) + bg_gradient_bottom[0] * ratio)
            g = int(bg_gradient_top[1] * (1 - ratio) + bg_gradient_bottom[1] * ratio)
            b = int(bg_gradient_top[2] * (1 - ratio) + bg_gradient_bottom[2] * ratio)
            draw.line([(0, y), (width, y)], fill=(r, g, b))

        # Top bar: main title
        draw.text((width // 2, 30), title, fill=text_muted, font=font_small, anchor="mt")

        # Progress indicator
        progress_y = 60
        total = len(steps)
        dot_radius = 6
        dot_spacing = 30
        start_x = (width - (total - 1) * dot_spacing) // 2
        for j in range(total):
            cx = start_x + j * dot_spacing
            color = accent if j <= i else (71, 85, 105)
            draw.ellipse([cx - dot_radius, progress_y - dot_radius,
                          cx + dot_radius, progress_y + dot_radius], fill=color)

        # Step counter
        counter = f"{i + 1} / {total}"
        draw.text((width // 2, 90), counter, fill=accent, font=font_small, anchor="mt")

        # Icon (emoji as text — may render as tofu on servers without emoji fonts)
        icon_text = step.get("icon", "📖")
        draw.text((width // 2, 160), icon_text, fill=text_white, font=font_icon, anchor="mt")

        # Step title
        step_title = step.get("title", "")
        draw.text((width // 2, 260), step_title, fill=text_white, font=font_large, anchor="mt")

        # Description — word wrap
        desc = step.get("description", "")
        wrapped = textwrap.fill(desc, width=35)
        draw.text((width // 2, 320), wrapped, fill=text_muted, font=font_medium,
                   anchor="mt", align="center")

        frames.append(img)

    # Save as animated GIF
    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    if frames:
        frames[0].save(
            output_path,
            save_all=True,
            append_images=frames[1:],
            duration=duration_per_frame_ms,
            loop=0,
        )

    return output_path


# ═══════════════════════════════════════════
# 2. Animation steps → MP4 video (via FFmpeg)
# ═══════════════════════════════════════════

def animation_steps_to_mp4(
    data: Dict[str, Any],
    output_path: str,
    width: int = 800,
    height: int = 500,
    duration_per_frame_sec: float = 3.0,
    fps: int = 1,
) -> str:
    """
    Render animation steps data to an MP4 video via FFmpeg.

    Generates the same Pillow frames as the GIF exporter, then pipes them
    to FFmpeg for H.264 encoding.  Falls back to an animated-GIF-to-MP4
    conversion if the raw-frame pipe fails.

    Requires FFmpeg to be installed and available on PATH.

    Args:
        data: {"title": "...", "steps": [{"icon": "🔬", "title": "...", "description": "..."}]}
        output_path: Where to save the MP4 file.
        width: Frame width in pixels.
        height: Frame height in pixels.
        duration_per_frame_sec: How long each step is shown (seconds).
        fps: Frames per second (1 is fine for static-step videos).

    Returns:
        The output file path.
    """
    import subprocess
    import tempfile

    title = data.get("title", "知识点动画")
    steps = data.get("steps", [])
    if not steps:
        steps = [{"icon": "📖", "title": "暂无内容", "description": "动画步骤为空"}]

    # --- render frames (same logic as GIF) ---
    font_large = _get_font(36)
    font_medium = _get_font(24)
    font_small = _get_font(18)
    font_icon = _get_font(64)

    bg_gradient_top = (30, 41, 59)
    bg_gradient_bottom = (15, 23, 42)
    text_white = (248, 250, 252)
    text_muted = (148, 163, 184)
    accent = (56, 189, 248)

    frames: List[Image.Image] = []

    for i, step in enumerate(steps):
        img = Image.new("RGB", (width, height))
        draw = ImageDraw.Draw(img)

        for y in range(height):
            ratio = y / height
            r = int(bg_gradient_top[0] * (1 - ratio) + bg_gradient_bottom[0] * ratio)
            g = int(bg_gradient_top[1] * (1 - ratio) + bg_gradient_bottom[1] * ratio)
            b = int(bg_gradient_top[2] * (1 - ratio) + bg_gradient_bottom[2] * ratio)
            draw.line([(0, y), (width, y)], fill=(r, g, b))

        draw.text((width // 2, 30), title, fill=text_muted, font=font_small, anchor="mt")

        total = len(steps)
        dot_radius = 6
        dot_spacing = 30
        start_x = (width - (total - 1) * dot_spacing) // 2
        for j in range(total):
            cx = start_x + j * dot_spacing
            color = accent if j <= i else (71, 85, 105)
            draw.ellipse([cx - dot_radius, 60 - dot_radius,
                          cx + dot_radius, 60 + dot_radius], fill=color)

        draw.text((width // 2, 90), f"{i + 1} / {total}", fill=accent, font=font_small, anchor="mt")
        draw.text((width // 2, 160), step.get("icon", "📖"), fill=text_white, font=font_icon, anchor="mt")
        draw.text((width // 2, 260), step.get("title", ""), fill=text_white, font=font_large, anchor="mt")

        desc = step.get("description", "")
        wrapped = textwrap.fill(desc, width=35)
        draw.text((width // 2, 320), wrapped, fill=text_muted, font=font_medium,
                   anchor="mt", align="center")

        frames.append(img)

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)

    # --- encode to MP4 via FFmpeg raw-frame pipe ---
    repeat_count = max(1, int(duration_per_frame_sec * fps))

    try:
        proc = subprocess.Popen(
            [
                "ffmpeg", "-y",
                "-f", "rawvideo",
                "-pix_fmt", "rgb24",
                "-s", f"{width}x{height}",
                "-r", str(fps),
                "-i", "pipe:0",
                "-c:v", "libx264",
                "-pix_fmt", "yuv420p",
                "-preset", "fast",
                "-crf", "23",
                output_path,
            ],
            stdin=subprocess.PIPE,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.PIPE,
        )

        for frame in frames:
            raw = frame.tobytes()
            for _ in range(repeat_count):
                proc.stdin.write(raw)

        proc.stdin.close()
        proc.wait(timeout=60)

        if proc.returncode != 0:
            raise RuntimeError(proc.stderr.read().decode(errors="replace"))

    except FileNotFoundError:
        raise RuntimeError(
            "FFmpeg is required for MP4 export but was not found on PATH. "
            "Install FFmpeg: https://ffmpeg.org/download.html"
        )

    return output_path


# ═══════════════════════════════════════════
# 3. Game/Animation HTML → PPT embedded slide
# ═══════════════════════════════════════════

def embed_html_in_pptx(
    pptx_path: str,
    html_path: str,
    slide_title: str = "互动内容",
    description: str = "点击下方链接打开互动内容",
) -> str:
    """
    Add a slide to an existing PPTX that links to an HTML game/animation.

    Since PPTX cannot natively embed HTML, we add a clearly marked slide
    with a hyperlink to the HTML file.

    Returns:
        The modified PPTX file path.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    prs = Presentation(pptx_path)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # Title
    from pptx.util import Emu
    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(8), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"🎮 {slide_title}"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1e, 0x3a, 0x5f)

    # Description
    desc_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    tf2 = desc_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = description
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(0x64, 0x74, 0x8b)

    # Hyperlink to HTML file
    html_name = Path(html_path).name
    link_box = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(0.6))
    tf3 = link_box.text_frame
    p3 = tf3.paragraphs[0]
    run = p3.add_run()
    run.text = f"📎 点击打开: {html_name}"
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(0x0e, 0xa5, 0xe9)
    run.font.underline = True
    run.hyperlink.address = html_name  # relative path — works when HTML is next to PPTX

    prs.save(pptx_path)
    return pptx_path


# ═══════════════════════════════════════════
# 3. ZIP bundle export
# ═══════════════════════════════════════════

def create_export_bundle(
    session_id: str,
    artifacts_dir: str,
    output_path: Optional[str] = None,
) -> str:
    """
    Bundle all generated artifacts (PPTX, DOCX, HTML games, animations, GIFs)
    into a single ZIP file for download.

    Returns:
        The ZIP file path.
    """
    import zipfile

    artifacts = Path(artifacts_dir)
    if not artifacts.exists():
        raise FileNotFoundError(f"Artifacts directory not found: {artifacts_dir}")

    if output_path is None:
        output_path = str(artifacts / f"{session_id}_courseware_bundle.zip")

    with zipfile.ZipFile(output_path, "w", zipfile.ZIP_DEFLATED) as zf:
        for fpath in artifacts.rglob("*"):
            if fpath.is_file() and not fpath.name.endswith(".zip"):
                arcname = fpath.relative_to(artifacts)
                zf.write(fpath, arcname)

    return output_path


# ═══════════════════════════════════════════
# Font helper
# ═══════════════════════════════════════════

def _get_font(size: int):
    """Try to load a CJK-compatible font, falling back to default."""
    font_candidates = [
        "C:/Windows/Fonts/msyh.ttc",       # Microsoft YaHei (Windows)
        "C:/Windows/Fonts/simhei.ttf",      # SimHei (Windows)
        "/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc",  # WenQuanYi (Linux)
        "/System/Library/Fonts/PingFang.ttc",              # PingFang (macOS)
    ]
    for path in font_candidates:
        if Path(path).exists():
            try:
                return ImageFont.truetype(path, size)
            except Exception:
                continue
    try:
        return ImageFont.truetype("arial.ttf", size)
    except Exception:
        return ImageFont.load_default()

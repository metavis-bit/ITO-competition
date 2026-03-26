"""
PPT 课件生成引擎

对应 A04 要求:
  - 4a) 根据指令集，自动生成包含封面、目录、内容页、总结页的PPT课件
  - 5c) 支持将最终课件以 .pptx 格式下载
"""

from __future__ import annotations

import json
import re
import time
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor


# ═══════════════════════════════════════════
# 颜色/字体常量
# ═══════════════════════════════════════════
_BLUE     = RGBColor(0x1A, 0x73, 0xE8)
_DARK     = RGBColor(0x20, 0x2A, 0x44)
_GRAY     = RGBColor(0x5F, 0x63, 0x68)
_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
_ACCENT   = RGBColor(0x34, 0xA8, 0x53)
_BG_LIGHT = RGBColor(0xF1, 0xF3, 0xF4)

_FONT_CN  = "微软雅黑"
_FONT_EN  = "Calibri"


# ═══════════════════════════════════════════
# 辅助函数
# ═══════════════════════════════════════════

def _set_text(shape, text: str, font_size: int = 18, bold: bool = False,
              color: RGBColor = _DARK, alignment=PP_ALIGN.LEFT,
              font_name: str = _FONT_CN):
    """设置文本框内容和样式"""
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name


def _add_textbox(slide, left, top, width, height, text: str,
                 font_size: int = 16, bold: bool = False,
                 color: RGBColor = _DARK, alignment=PP_ALIGN.LEFT):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    _set_text(txBox, text, font_size, bold, color, alignment)
    return txBox


def _add_bullet_textbox(slide, left, top, width, height, items: List[str],
                        font_size: int = 16, color: RGBColor = _DARK):
    """添加带项目符号的文本框"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(6)
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = f"• {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = _FONT_CN
    return txBox


def _add_shape_rect(slide, left, top, width, height, fill_color: RGBColor):
    """添加矩形色块"""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


# ═══════════════════════════════════════════
# 幻灯片生成器
# ═══════════════════════════════════════════

def _make_cover_slide(prs: Presentation, title: str, subtitle: str, author: str = ""):
    """封面页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    # 蓝色背景条
    _add_shape_rect(slide, 0, 0, 10, 2.5, _BLUE)
    # 标题
    _add_textbox(slide, 0.8, 0.6, 8.4, 1.2, title, font_size=36, bold=True, color=_WHITE, alignment=PP_ALIGN.CENTER)
    # 副标题
    _add_textbox(slide, 0.8, 1.8, 8.4, 0.6, subtitle, font_size=18, color=_WHITE, alignment=PP_ALIGN.CENTER)
    # 作者
    if author:
        _add_textbox(slide, 0.8, 3.2, 8.4, 0.5, author, font_size=14, color=_GRAY, alignment=PP_ALIGN.CENTER)
    # 日期
    _add_textbox(slide, 0.8, 3.8, 8.4, 0.5, datetime.now().strftime("%Y年%m月%d日"),
                 font_size=12, color=_GRAY, alignment=PP_ALIGN.CENTER)
    return slide


def _make_toc_slide(prs: Presentation, sections: List[str]):
    """目录页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, 0.5, 0.3, 9, 0.8, "目  录", font_size=32, bold=True, color=_BLUE, alignment=PP_ALIGN.CENTER)
    # 目录分隔线
    _add_shape_rect(slide, 2, 1.1, 6, 0.03, _BLUE)

    for i, section in enumerate(sections):
        y = 1.5 + i * 0.6
        # 序号圆
        _add_textbox(slide, 1.2, y, 0.5, 0.5, f"{i+1:02d}", font_size=14, bold=True, color=_BLUE)
        # 标题
        _add_textbox(slide, 2.0, y, 6.5, 0.5, section, font_size=18, color=_DARK)
    return slide


def _make_section_title_slide(prs: Presentation, section_num: int, section_title: str):
    """章节标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_shape_rect(slide, 0, 1.5, 10, 2.5, _BLUE)
    _add_textbox(slide, 0.8, 1.0, 2, 0.6, f"PART {section_num:02d}", font_size=16, color=_BLUE, alignment=PP_ALIGN.LEFT)
    _add_textbox(slide, 0.8, 2.0, 8.4, 1.0, section_title, font_size=32, bold=True, color=_WHITE, alignment=PP_ALIGN.CENTER)
    return slide


def _add_info_card(slide, left, top, width, height, title_text: str,
                   body_text: str, accent_color: RGBColor = _BLUE):
    """添加带色块标题的信息卡片（图文并茂的视觉元素）"""
    from pptx.enum.shapes import MSO_SHAPE
    # 卡片背景
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(left), Inches(top), Inches(width), Inches(height))
    bg.fill.solid()
    bg.fill.fore_color.rgb = _BG_LIGHT
    bg.line.fill.background()
    # 色条
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(left), Inches(top), Inches(0.08), Inches(height))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()
    # 卡片标题
    _add_textbox(slide, left + 0.2, top + 0.1, width - 0.3, 0.35,
                 title_text, font_size=13, bold=True, color=accent_color)
    # 卡片正文
    _add_textbox(slide, left + 0.2, top + 0.45, width - 0.3, height - 0.55,
                 body_text, font_size=11, color=_DARK)


def _add_step_circle(slide, left, top, number: int, label: str, color: RGBColor = _BLUE):
    """添加步骤圆圈（用于流程图式布局）"""
    from pptx.enum.shapes import MSO_SHAPE
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                    Inches(left), Inches(top), Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    # 数字
    _add_textbox(slide, left + 0.05, top + 0.08, 0.5, 0.4,
                 str(number), font_size=16, bold=True, color=_WHITE, alignment=PP_ALIGN.CENTER)
    # 标签
    _add_textbox(slide, left - 0.15, top + 0.65, 0.9, 0.4,
                 label, font_size=9, color=_DARK, alignment=PP_ALIGN.CENTER)


def _make_content_slide(prs: Presentation, title: str, points: List[str],
                        speaker_notes: str = "", visual_hint: str = "",
                        layout: str = "auto", image_path: str = ""):
    """
    内容页（图文并茂版）

    layout: "auto"(自动选择), "bullets"(要点列表), "cards"(信息卡片),
            "two_col"(双栏), "steps"(步骤流程)
    image_path: 可选，插入到右侧的知识图片路径
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 顶部蓝色条
    _add_shape_rect(slide, 0, 0, 10, 0.06, _BLUE)
    # 标题
    _add_textbox(slide, 0.5, 0.2, 9, 0.6, title, font_size=24, bold=True, color=_DARK)
    # 分隔线
    _add_shape_rect(slide, 0.5, 0.85, 9, 0.02, _BG_LIGHT)

    # 自动选择布局
    if layout == "auto":
        if len(points) <= 3 and all(len(p) > 30 for p in points):
            layout = "cards"
        elif len(points) >= 4 and len(points) <= 6 and all(len(p) < 25 for p in points):
            layout = "steps"
        elif len(points) > 3:
            layout = "two_col"
        else:
            layout = "bullets"

    if layout == "cards" and points:
        # 信息卡片布局：每个要点一个卡片
        colors = [_BLUE, _ACCENT, RGBColor(0xEA, 0x43, 0x35), RGBColor(0xFB, 0xBC, 0x04)]
        n = min(len(points), 4)
        card_w = (9.0 - 0.2 * (n - 1)) / n
        for i, pt in enumerate(points[:4]):
            parts = pt.split("：", 1) if "：" in pt else pt.split(":", 1) if ":" in pt else [f"要点{i+1}", pt]
            card_title = parts[0].strip()
            card_body = parts[1].strip() if len(parts) > 1 else pt
            x = 0.5 + i * (card_w + 0.2)
            _add_info_card(slide, x, 1.1, card_w, 2.8,
                           card_title, card_body, colors[i % len(colors)])

    elif layout == "steps" and points:
        # 步骤流程布局：圆圈 + 箭头
        n = min(len(points), 6)
        step_w = 8.0 / n
        for i, pt in enumerate(points[:6]):
            x = 1.0 + i * step_w
            _add_step_circle(slide, x, 1.3, i + 1, pt[:15], _BLUE if i % 2 == 0 else _ACCENT)
            # 箭头
            if i < n - 1:
                _add_textbox(slide, x + 0.7, 1.4, 0.4, 0.3, "→", font_size=18, color=_GRAY,
                             alignment=PP_ALIGN.CENTER)
        # 详细说明
        if any(len(p) > 15 for p in points):
            detail_text = "\n".join(f"{i+1}. {p}" for i, p in enumerate(points))
            _add_textbox(slide, 0.7, 2.6, 8.6, 2.5, detail_text, font_size=12, color=_DARK)

    elif layout == "two_col" and points:
        # 双栏布局：左文字，右视觉区域
        mid = len(points) // 2
        left_pts = points[:mid] if mid > 0 else points[:2]
        right_pts = points[mid:] if mid > 0 else points[2:]

        _add_bullet_textbox(slide, 0.5, 1.1, 4.3, 3.8, left_pts, font_size=14, color=_DARK)
        # 右侧用信息卡片
        for i, pt in enumerate(right_pts[:3]):
            _add_info_card(slide, 5.2, 1.1 + i * 1.35, 4.3, 1.2,
                           f"要点 {mid + i + 1}", pt,
                           [_BLUE, _ACCENT, RGBColor(0xEA, 0x43, 0x35)][i % 3])
    else:
        # 默认要点列表（带左侧装饰条）
        _add_shape_rect(slide, 0.5, 1.1, 0.06, 3.8, _BLUE)
        if points:
            _add_bullet_textbox(slide, 0.8, 1.1, 8.5, 3.8, points, font_size=15, color=_DARK)

    # 视觉提示
    if visual_hint:
        _add_textbox(slide, 5.5, 4.9, 4.2, 0.4, f"📌 {visual_hint}", font_size=9, color=_GRAY)

    # ★ 插入知识图片（图文并茂的关键）
    if image_path:
        from pathlib import Path as _P
        if _P(image_path).exists():
            try:
                from PIL import Image as _Img
                img = _Img.open(image_path)
                w, h = img.size
                aspect = h / w

                # 根据布局决定图片位置和大小
                if layout in ("bullets", "auto"):
                    # 右侧放图：缩小文字区域到左边5寸，图片放右边4寸
                    img_w_in = 3.8
                    img_h_in = min(img_w_in * aspect, 3.5)
                    img_left = Inches(5.8)
                    img_top = Inches(1.1)
                    # 在图片下方加标签
                    slide.shapes.add_picture(image_path, img_left, img_top,
                                             Inches(img_w_in), Inches(img_h_in))
                    _add_textbox(slide, 5.8, 1.1 + img_h_in + 0.05, 3.8, 0.3,
                                 "📷 知识图解", font_size=9, color=_GRAY,
                                 alignment=PP_ALIGN.CENTER)
                elif layout == "two_col":
                    # 双栏模式：图片替换右栏
                    img_w_in = 4.0
                    img_h_in = min(img_w_in * aspect, 3.5)
                    slide.shapes.add_picture(image_path,
                                             Inches(5.5), Inches(1.1),
                                             Inches(img_w_in), Inches(img_h_in))
                else:
                    # cards/steps等：底部小图
                    img_w_in = 2.5
                    img_h_in = min(img_w_in * aspect, 1.8)
                    slide.shapes.add_picture(image_path,
                                             Inches(7.0), Inches(3.5),
                                             Inches(img_w_in), Inches(img_h_in))
            except Exception:
                pass  # 图片插入失败不影响PPT

    # 演讲备注
    if speaker_notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = speaker_notes

    return slide


def _make_summary_slide(prs: Presentation, title: str, key_points: List[str]):
    """总结页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_shape_rect(slide, 0, 0, 10, 1.2, _ACCENT)
    _add_textbox(slide, 0.5, 0.3, 9, 0.7, f"📌 {title}", font_size=28, bold=True, color=_WHITE, alignment=PP_ALIGN.CENTER)

    if key_points:
        _add_bullet_textbox(slide, 1.0, 1.6, 8, 3.5, key_points, font_size=18, color=_DARK)
    return slide


def _make_thanks_slide(prs: Presentation):
    """结束页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_shape_rect(slide, 0, 2, 10, 2, _BLUE)
    _add_textbox(slide, 0, 2.3, 10, 1, "谢谢聆听", font_size=40, bold=True, color=_WHITE, alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, 0, 3.3, 10, 0.6, "欢迎提问与交流", font_size=18, color=_WHITE, alignment=PP_ALIGN.CENTER)
    return slide


def _make_interactive_slide(prs: Presentation, game_title: str, game_filename: str,
                            game_type_desc: str = "HTML5互动游戏"):
    """
    互动环节页 — 展示游戏信息 + 提供HTML5文件名提示。

    对应 A04 5c): "动画或小游戏可以以html5网页导出或集成到PPT中"
    教师可在演示时打开配套的HTML游戏文件。
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 背景色块
    _add_shape_rect(slide, 0, 0, 10, 5.625, RGBColor(0x0D, 0x11, 0x2B))
    # 装饰条
    _add_shape_rect(slide, 0, 0, 10, 0.06, _ACCENT)

    # 标题
    _add_textbox(slide, 0, 0.5, 10, 0.8, "🎮 课堂互动环节",
                 font_size=32, bold=True, color=_ACCENT, alignment=PP_ALIGN.CENTER)

    # 游戏信息卡片
    from pptx.enum.shapes import MSO_SHAPE
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(2), Inches(1.6), Inches(6), Inches(2.5))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
    card.line.color.rgb = RGBColor(0x38, 0xBD, 0xF8)

    _add_textbox(slide, 2.3, 1.9, 5.4, 0.6, f"📋 {game_title}",
                 font_size=22, bold=True, color=_WHITE, alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, 2.3, 2.5, 5.4, 0.5, f"类型: {game_type_desc}",
                 font_size=14, color=RGBColor(0x94, 0xA3, 0xB8), alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, 2.3, 3.1, 5.4, 0.6, f"📁 请打开配套文件: {game_filename}",
                 font_size=13, color=RGBColor(0x38, 0xBD, 0xF8), alignment=PP_ALIGN.CENTER)

    # 操作提示
    _add_textbox(slide, 1, 4.5, 8, 0.5,
                 "💡 在浏览器中打开上述HTML文件，即可开始课堂互动",
                 font_size=12, color=RGBColor(0x94, 0xA3, 0xB8), alignment=PP_ALIGN.CENTER)
    return slide


# ═══════════════════════════════════════════
# PPT 生成引擎
# ═══════════════════════════════════════════

class PPTGenerator:
    """
    PPT 课件自动生成器。

    可以两种方式驱动:
    1. 传入结构化指令 (ppt_outline JSON) → 直接生成
    2. 传入自然语言需求 → LLM 生成大纲 → 再生成PPT

    Example:
        gen = PPTGenerator(config_path="config.yaml")
        result = gen.generate_from_outline(outline_data, output_dir="outputs/pptx")
        print(result["pptx_path"])
    """

    def __init__(self, config_path: str = "config.yaml"):
        self.config_path = config_path
        self._generator = None

    @property
    def generator(self):
        if self._generator is None:
            from .config import load_config
            from .generator import GeneratorConfig, QwenGenerator
            cfg = load_config(self.config_path)
            g_cfg = cfg.get("generator", {}) or {}
            self._generator = QwenGenerator(GeneratorConfig(
                mode=str(g_cfg.get("mode", "api")),
                api_model=str(g_cfg.get("api_model", "qwen-plus-latest")),
                api_base_url_env=str(g_cfg.get("api_base_url_env", "OPENAI_BASE_URL")),
                api_key_env=str(g_cfg.get("api_key_env", "OPENAI_API_KEY")),
                timeout_sec=float(g_cfg.get("timeout_sec", 120.0)),
                local_model_path=str(g_cfg.get("local_model_path", "")),
                device=str(g_cfg.get("device", "cpu")),
                max_new_tokens=int(g_cfg.get("max_new_tokens", 4096)),
                temperature=float(g_cfg.get("temperature", 0.3)),
            ))
        return self._generator

    def generate_from_outline(
        self,
        outline: Dict[str, Any],
        output_dir: str = "outputs/pptx",
        author: str = "",
    ) -> Dict[str, Any]:
        """
        根据结构化大纲生成PPT。

        outline 格式:
        {
          "title": "课件标题",
          "subtitle": "副标题",
          "slides": [
            {
              "type": "cover/toc/section_title/content/summary/thanks",
              "title": "幻灯片标题",
              "content_points": ["要点1", "要点2"],
              "speaker_notes": "演讲备注",
              "visual_suggestion": "视觉建议",
              "section_num": 1  // for section_title type
            }
          ]
        }
        """
        t0 = time.time()
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)  # 16:9

        main_title = outline.get("title", "教学课件")
        subtitle = outline.get("subtitle", "")
        slides_data = outline.get("slides", [])

        # 如果没有明确的slides结构，自动组装
        if not slides_data:
            slides_data = self._auto_structure(outline)

        # 自动图片匹配
        used_images = set()
        _find_image = None
        try:
            from .image_matcher import find_best_image
            _find_image = find_best_image
        except Exception:
            pass

        for sd in slides_data:
            stype = sd.get("type", "content")
            title = sd.get("title", "")
            points = sd.get("content_points", [])
            notes = sd.get("speaker_notes", "")
            visual = sd.get("visual_suggestion", "")

            if stype == "cover":
                _make_cover_slide(prs, title or main_title, subtitle or sd.get("subtitle", ""), author)
            elif stype == "toc":
                sections = sd.get("sections", points)
                _make_toc_slide(prs, sections)
            elif stype == "section_title":
                _make_section_title_slide(prs, sd.get("section_num", 1), title)
            elif stype == "summary":
                _make_summary_slide(prs, title or "本课总结", points)
            elif stype == "thanks":
                _make_thanks_slide(prs)
            else:  # content
                layout = sd.get("layout", "auto")
                # 自动匹配图片
                img_path = ""
                if _find_image:
                    matched = _find_image(title, points, used_images=used_images)
                    if matched:
                        img_path = matched
                        from pathlib import Path as _PP
                        used_images.add(_PP(matched).name)
                _make_content_slide(prs, title, points, notes, visual, layout=layout, image_path=img_path)

        # 保存
        Path(output_dir).mkdir(parents=True, exist_ok=True)
        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe_title = re.sub(r'[^\w\u4e00-\u9fff]', '_', main_title)[:30]
        filename = f"ppt_{safe_title}_{ts}.pptx"
        filepath = str(Path(output_dir) / filename)
        prs.save(filepath)

        gen_time = round(time.time() - t0, 2)
        return {
            "pptx_path": filepath,
            "title": main_title,
            "slide_count": len(prs.slides),
            "generation_time": gen_time,
        }

    def generate_from_text(
        self,
        topic: str,
        teaching_goal: str = "",
        key_points: str = "",
        style: str = "清晰简洁",
        slide_count: int = 10,
        context: str = "",
        output_dir: str = "outputs/pptx",
    ) -> Dict[str, Any]:
        """
        从自然语言需求生成PPT（先调LLM生成大纲，再渲染）。
        """
        prompt = f"""你是一位专业的教学课件设计师。请根据以下需求，设计一份PPT课件大纲。

【课题】{topic}
【教学目标】{teaching_goal or '无特殊要求'}
【重点内容】{key_points or '无特殊要求'}
【风格】{style}
【幻灯片数量】约{slide_count}页
【参考内容】
{context[:3000] if context else '无'}

请严格按以下JSON格式输出，不要输出多余文字：
{{
  "title": "课件标题",
  "subtitle": "副标题（学科+章节）",
  "slides": [
    {{"type": "cover", "title": "课件标题", "subtitle": "副标题"}},
    {{"type": "toc", "title": "目录", "sections": ["第一部分标题", "第二部分标题", "..."]}},
    {{"type": "section_title", "title": "第一部分标题", "section_num": 1}},
    {{"type": "content", "title": "内容标题", "content_points": ["要点1", "要点2", "要点3"], "layout": "cards", "speaker_notes": "演讲备注", "visual_suggestion": "图表建议"}},
    {{"type": "content", "title": "流程步骤", "content_points": ["步骤1", "步骤2", "步骤3", "步骤4"], "layout": "steps"}},
    {{"type": "content", "title": "对比分析", "content_points": ["左侧概念A", "左侧概念B", "右侧概念C", "右侧概念D"], "layout": "two_col"}},
    {{"type": "summary", "title": "本课总结", "content_points": ["核心要点1", "核心要点2"]}},
    {{"type": "thanks"}}
  ]
}}

layout 可选值说明：
- "cards": 每个要点一张信息卡片（适合3个以内的核心概念解释）
- "steps": 步骤流程圆圈（适合有顺序的过程/步骤）
- "two_col": 双栏对比（适合4个以上要点或对比分析）
- "bullets": 普通要点列表（默认）"""

        raw = self.generator.generate(prompt)
        outline = self._parse_json(raw)

        if not outline.get("slides"):
            outline = {
                "title": topic,
                "subtitle": teaching_goal,
                "slides": [
                    {"type": "cover", "title": topic, "subtitle": teaching_goal},
                    {"type": "content", "title": topic, "content_points": [key_points or "内容待补充"]},
                    {"type": "thanks"},
                ],
            }

        return self.generate_from_outline(outline, output_dir)

    def _auto_structure(self, outline: Dict) -> List[Dict]:
        """当outline没有slides字段时，从其他字段自动组装"""
        slides = []
        title = outline.get("title", "课件")
        slides.append({"type": "cover", "title": title, "subtitle": outline.get("subtitle", "")})

        # 从 content_blocks / knowledge_points 等构建内容页
        blocks = outline.get("content_blocks", [])
        kps = outline.get("knowledge_points", [])

        if blocks:
            sections = [b.get("title", f"第{i+1}部分") for i, b in enumerate(blocks)]
            slides.append({"type": "toc", "sections": sections})
            for i, block in enumerate(blocks):
                slides.append({"type": "section_title", "title": block.get("title", ""), "section_num": i + 1})
                slides.append({
                    "type": "content",
                    "title": block.get("title", ""),
                    "content_points": [block.get("content_hint", "内容要点")],
                })
        elif kps:
            sections = [kp.get("name", f"知识点{i+1}") for i, kp in enumerate(kps)]
            slides.append({"type": "toc", "sections": sections})
            for i, kp in enumerate(kps):
                slides.append({
                    "type": "content",
                    "title": kp.get("name", ""),
                    "content_points": [kp.get("description", ""), f"重要程度: {kp.get('importance', '')}"],
                })

        # 总结
        key_focus = outline.get("key_focus", [])
        if key_focus:
            slides.append({"type": "summary", "title": "本课总结", "content_points": key_focus})

        slides.append({"type": "thanks"})
        return slides

    @staticmethod
    def _parse_json(raw: str) -> Dict[str, Any]:
        text = raw.strip()
        text = re.sub(r'^```(?:json)?\s*', '', text, flags=re.MULTILINE)
        text = re.sub(r'```\s*$', '', text, flags=re.MULTILINE)
        text = text.strip()
        try:
            return json.loads(text)
        except json.JSONDecodeError:
            pass
        match = re.search(r'\{[\s\S]*\}', text)
        if match:
            candidate = re.sub(r',\s*([}\]])', r'\1', match.group(0))
            try:
                return json.loads(candidate)
            except json.JSONDecodeError:
                pass
        return {}
